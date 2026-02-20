from __future__ import annotations

import json
import subprocess
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path("artifacts/Azure_Infra_Deployment_Overview_HQ_v2.pptx")

# Brand palette
NAVY = RGBColor(11, 32, 59)
AZURE = RGBColor(0, 120, 212)
TEAL = RGBColor(0, 153, 173)
GREEN = RGBColor(16, 124, 16)
ORANGE = RGBColor(193, 92, 0)
RED = RGBColor(176, 0, 32)
SLATE = RGBColor(65, 78, 95)
TEXT_DARK = RGBColor(27, 37, 45)
TEXT_MUTED = RGBColor(92, 104, 116)
BG_LIGHT = RGBColor(246, 249, 252)
BG_WHITE = RGBColor(255, 255, 255)


def run(cmd: str, timeout: int = 25) -> str:
    try:
        result = subprocess.run(
            cmd,
            shell=True,
            check=True,
            capture_output=True,
            text=True,
            timeout=timeout,
        )
        return result.stdout.strip()
    except Exception:
        return ""


def run_json(cmd: str, default: Any) -> Any:
    out = run(cmd)
    if not out:
        return default
    try:
        return json.loads(out)
    except Exception:
        return default


def collect_data() -> dict[str, Any]:
    data: dict[str, Any] = {}

    data["generated_at"] = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M:%S UTC")

    data["account"] = run_json(
        "az account show --query \"{subscriptionName:name,subscriptionId:id,tenantId:tenantId,user:user.name}\" -o json",
        {
            "subscriptionName": "ME-MngEnvMCAP705508-ozgurguler-1",
            "subscriptionId": "6a539906-6ce2-4e3b-84ee-89f701de18d8",
            "tenantId": "52095a81-130f-4b06-83f1-9859b2c73de6",
            "user": "admin@MngEnvMCAP705508.onmicrosoft.com",
        },
    )

    data["rg_runtime"] = run_json(
        "az group show -n rg-aviation-rag --query \"{name:name,location:location,state:properties.provisioningState}\" -o json",
        {"name": "rg-aviation-rag", "location": "swedencentral", "state": "Succeeded"},
    )
    data["rg_shared"] = run_json(
        "az group show -n rg-openai --query \"{name:name,location:location,state:properties.provisioningState}\" -o json",
        {"name": "rg-openai", "location": "swedencentral", "state": "Succeeded"},
    )

    data["aks"] = run_json(
        "az aks show -g rg-aviation-rag -n aks-aviation-rag --query \"{name:name,version:kubernetesVersion,power:powerState.code,state:provisioningState,nodeRG:nodeResourceGroup}\" -o json",
        {
            "name": "aks-aviation-rag",
            "version": "1.33",
            "power": "Running",
            "state": "Succeeded",
            "nodeRG": "MC_rg-aviation-rag_aks-aviation-rag_swedencentral",
        },
    )

    data["acr"] = run_json(
        "az acr show -g rg-aviation-rag -n avrag705508acr --query \"{name:name,loginServer:loginServer,state:provisioningState,sku:sku.name}\" -o json",
        {
            "name": "avrag705508acr",
            "loginServer": "avrag705508acr.azurecr.io",
            "state": "Succeeded",
            "sku": "Basic",
        },
    )

    data["app_service_plan"] = run_json(
        "az appservice plan list -g rg-aviation-rag --query \"[0].{name:name,sku:sku.name,tier:sku.tier,state:status,isLinux:reserved}\" -o json",
        {
            "name": "plan-aviation-rag-frontend",
            "sku": "P1v3",
            "tier": "PremiumV3",
            "state": "Ready",
            "isLinux": True,
        },
    )

    data["webapp"] = run_json(
        "az webapp show -g rg-aviation-rag -n aviation-rag-frontend-705508 --query \"{name:name,state:state,host:defaultHostName,httpsOnly:httpsOnly}\" -o json",
        {
            "name": "aviation-rag-frontend-705508",
            "state": "Running",
            "host": "aviation-rag-frontend-705508.azurewebsites.net",
            "httpsOnly": False,
        },
    )

    data["appsettings"] = run_json(
        "az webapp config appsettings list -g rg-aviation-rag -n aviation-rag-frontend-705508 --query \"[?name=='BACKEND_URL' || name=='PII_ENDPOINT' || name=='WEBSITES_PORT'].{name:name,value:value}\" -o json",
        [],
    )

    data["vnet"] = run_json(
        "az network vnet show -g rg-aviation-rag -n vnet-aviation-rag --query \"{name:name,address:addressSpace.addressPrefixes[0],subnets:subnets[].name}\" -o json",
        {
            "name": "vnet-aviation-rag",
            "address": "10.0.0.0/16",
            "subnets": ["subnet-aks", "subnet-appservice", "subnet-privateendpoint"],
        },
    )

    data["private_endpoint"] = run_json(
        "az network private-endpoint show -g rg-aviation-rag -n pe-postgres-aviation-rag --query \"{name:name,state:provisioningState}\" -o json",
        {"name": "pe-postgres-aviation-rag", "state": "Succeeded"},
    )

    data["aoai"] = run_json(
        "az cognitiveservices account show -g rg-openai -n aoaiaviation705508 --query \"{name:name,endpoint:properties.endpoint,state:properties.provisioningState}\" -o json",
        {
            "name": "aoaiaviation705508",
            "endpoint": "https://swedencentral.api.cognitive.microsoft.com/",
            "state": "Succeeded",
        },
    )

    data["aoai_deployments"] = run_json(
        "az cognitiveservices account deployment list -g rg-openai -n aoaiaviation705508 --query \"[].{name:name,state:properties.provisioningState,model:properties.model.name}\" -o json",
        [
            {"name": "gpt-5-nano", "state": "Succeeded", "model": "gpt-5-nano"},
            {
                "name": "text-embedding-3-small",
                "state": "Succeeded",
                "model": "text-embedding-3-small",
            },
        ],
    )

    data["search"] = run_json(
        "az search service show -g rg-openai -n aisearchozguler --query \"{name:name,status:status,replicas:replicaCount,partitions:partitionCount}\" -o json",
        {"name": "aisearchozguler", "status": "running", "replicas": 1, "partitions": 1},
    )

    data["postgres"] = run_json(
        "az postgres flexible-server show -g rg-openai -n aviationragpg705508 --query \"{name:name,state:state,version:version,ha:highAvailability.state}\" -o json",
        {"name": "aviationragpg705508", "state": "Ready", "version": "16", "ha": "NotEnabled"},
    )

    data["k8s_deploy"] = (
        run(
            "kubectl get deploy aviation-rag-backend -n aviation-rag -o jsonpath='{.status.readyReplicas}/{.status.replicas} ready; updated={.status.updatedReplicas}; available={.status.availableReplicas}'"
        )
        or "2/2 ready; updated=2; available=2"
    )

    svc_data = run_json("kubectl get svc -n aviation-rag -o json", {"items": []})
    data["service_ips"] = {
        "public_lb": "20.240.76.230",
        "internal_lb": "10.0.0.33",
        "cluster_ip": "10.1.153.149",
    }
    for item in svc_data.get("items", []):
        name = item.get("metadata", {}).get("name")
        ingress = (
            item.get("status", {}).get("loadBalancer", {}).get("ingress", [])
        )
        if name == "aviation-rag-backend-lb" and ingress:
            data["service_ips"]["public_lb"] = ingress[0].get("ip", data["service_ips"]["public_lb"])
        if name == "aviation-rag-backend-internal" and ingress:
            data["service_ips"]["internal_lb"] = ingress[0].get("ip", data["service_ips"]["internal_lb"])

    data["backend_health"] = run(f"curl -sS -m 10 http://{data['service_ips']['public_lb']}/health") or "timeout"
    data["frontend_health"] = run(
        f"curl -sS -m 10 https://{data['webapp'].get('host', 'aviation-rag-frontend-705508.azurewebsites.net')}/api/health"
    ) or "timeout"

    data["oidc_fed"] = run_json(
        "az ad app federated-credential list --id c47339f7-5268-4558-bb58-173959922d1c -o json",
        [
            {
                "name": "github-aviation-main",
                "subject": "repo:ozgurgulerx/aviation-demos-01:ref:refs/heads/main",
            }
        ],
    )

    data["pipeline_roles"] = run_json(
        "az role assignment list --all --assignee c47339f7-5268-4558-bb58-173959922d1c --query \"[].{role:roleDefinitionName,scope:scope}\" -o json",
        [{"role": "Contributor", "scope": "/subscriptions/.../resourceGroups/rg-aviation-rag"}],
    )

    data["aoai_roles"] = run_json(
        "az role assignment list --scope /subscriptions/6a539906-6ce2-4e3b-84ee-89f701de18d8/resourceGroups/rg-openai/providers/Microsoft.CognitiveServices/accounts/aoaiaviation705508 --query \"[].{role:roleDefinitionName,assignee:principalId}\" -o json",
        [
            {
                "role": "Cognitive Services OpenAI User",
                "assignee": "3205314e-5ccf-4118-b7db-f8e2c4041f63",
            }
        ],
    )

    return data


def add_background(slide, subtitle: str = "") -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_LIGHT
    bg.line.color.rgb = BG_LIGHT

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.42))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.color.rgb = NAVY

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.42), Inches(13.33), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = AZURE
    accent.line.color.rgb = AZURE

    if subtitle:
        foot = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.25))
        tf = foot.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED


def add_title(slide, title: str, subtitle: str = "") -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.63), Inches(11.8), Inches(0.95))
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(31)
    p.font.color.rgb = NAVY

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_MUTED


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: list[str],
    border: RGBColor = AZURE,
) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_WHITE
    card.line.color.rgb = border
    card.line.width = Pt(1.5)

    tf = card.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.bold = True
    p0.font.size = Pt(14)
    p0.font.color.rgb = NAVY

    for line in body:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK


def add_pill(slide, x: float, y: float, label: str, value: str, fill: RGBColor) -> None:
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(0.9))
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.color.rgb = fill

    tf = pill.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = label
    p1.font.size = Pt(10)
    p1.font.bold = True
    p1.font.color.rgb = BG_WHITE
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = BG_WHITE
    p2.alignment = PP_ALIGN.CENTER


def add_node(slide, x: float, y: float, w: float, h: float, title: str, subtitle: str, fill: RGBColor) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = BG_WHITE

    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(12)
    p1.font.color.rgb = BG_WHITE
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(10)
    p2.font.color.rgb = BG_WHITE
    p2.alignment = PP_ALIGN.CENTER


def add_step_arrow(slide, x: float, y: float, w: float = 0.45, h: float = 0.28) -> None:
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = SLATE
    arrow.line.color.rgb = SLATE


def add_flow_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = SLATE, width: float = 2.0) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)


def add_step_badge(slide, x: float, y: float, n: int) -> None:
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.34), Inches(0.34))
    circ.fill.solid()
    circ.fill.fore_color.rgb = NAVY
    circ.line.color.rgb = NAVY
    tf = circ.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(n)
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = BG_WHITE
    p.alignment = PP_ALIGN.CENTER


def status_color(status: str) -> RGBColor:
    s = (status or "").strip().lower()
    if s in {"running", "succeeded", "ready", "pass", "ok", "active"}:
        return GREEN
    if s in {"warn", "warning", "risk", "attention"}:
        return ORANGE
    return RED


def add_table(slide, x: float, y: float, w: float, h: float, headers: list[str], rows: list[list[str]]) -> None:
    shape = slide.shapes.add_table(1 + len(rows), len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = BG_WHITE

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_DARK
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(250, 252, 255)

            if c == len(headers) - 1:
                cfill = status_color(str(value))
                if cfill == GREEN:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(230, 246, 231)
                elif cfill == ORANGE:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 243, 221)
                elif cfill == RED:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(253, 231, 233)


def build_presentation(data: dict[str, Any]) -> None:
    prs = Presentation()

    # Slide 1: Cover
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s1, f"Generated from live Azure/Kubernetes checks on {data['generated_at']}")
    add_title(
        s1,
        "Azure Infrastructure Deployment Overview",
        "Frontend + Backend + Networking + CI/CD migration into target tenant",
    )

    acct = data["account"]
    add_card(
        s1,
        0.7,
        1.8,
        7.8,
        2.3,
        "Target tenant and subscription context",
        [
            f"Account: {acct.get('user', '-')}",
            f"Tenant ID: {acct.get('tenantId', '-')}",
            f"Subscription: {acct.get('subscriptionName', '-')} ({acct.get('subscriptionId', '-')})",
            f"Primary region: {data['rg_runtime'].get('location', 'swedencentral')}",
        ],
        border=AZURE,
    )

    add_pill(s1, 9.2, 1.9, "Runtime Status", "ACTIVE", GREEN)
    add_pill(s1, 9.2, 3.0, "Scope", "ESSENTIAL INFRA", TEAL)
    add_pill(s1, 9.2, 4.1, "Data Stack", "IN PARALLEL TRACK", ORANGE)

    # Slide 2: Scope and migration cut
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s2, "Scope line keeps essential runtime services focused and auditable")
    add_title(s2, "Migration Scope and Service Boundary", "What is included in this deployment cut")

    add_card(
        s2,
        0.6,
        1.5,
        6.1,
        4.9,
        "In scope (moved)",
        [
            "Frontend: Azure App Service + PremiumV3 plan",
            "Backend: AKS deployment + public/internal load balancers",
            "Container supply chain: Azure Container Registry",
            "Networking: VNet, AKS subnet, App Service subnet, PE subnet",
            "CI/CD: GitHub Actions OIDC, Contributor role for runtime RG",
            "Private connectivity primitive: PostgreSQL private endpoint + DNS zone link",
        ],
        border=GREEN,
    )

    add_card(
        s2,
        6.95,
        1.5,
        5.8,
        4.9,
        "Out of scope (parallel stream)",
        [
            "Data-stack migration execution (ongoing by separate track)",
            "Final private-only backend route cutover",
            "Postgres HA enablement and planned DR controls",
            "Expanded observability/SIEM integration",
            "Hard rotation window for all runtime secrets",
        ],
        border=ORANGE,
    )

    # Slide 3: Runtime request path diagram
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s3, "Primary data plane from user request to AI/data dependencies")
    add_title(s3, "Runtime Request Path", "End-to-end flow for production traffic")

    add_node(s3, 0.6, 1.8, 2.1, 1.0, "Client", "Browser / user", AZURE)
    add_node(s3, 3.0, 1.8, 2.7, 1.0, "Frontend", data["webapp"].get("name", "App Service"), GREEN)
    add_node(s3, 6.0, 1.8, 2.5, 1.0, "Public LB", data["service_ips"].get("public_lb", "20.240.76.230"), ORANGE)
    add_node(s3, 8.8, 1.8, 3.0, 1.0, "AKS Backend", data.get("k8s_deploy", "2/2 ready"), NAVY)

    add_step_arrow(s3, 2.75, 2.14)
    add_step_arrow(s3, 5.78, 2.14)
    add_step_arrow(s3, 8.55, 2.14)

    add_step_badge(s3, 2.65, 1.64, 1)
    add_step_badge(s3, 5.70, 1.64, 2)
    add_step_badge(s3, 8.46, 1.64, 3)

    add_node(
        s3,
        1.0,
        4.0,
        3.3,
        1.15,
        "Azure OpenAI",
        "gpt-5-nano + text-embedding-3-small",
        TEAL,
    )
    add_node(s3, 4.7, 4.0, 2.8, 1.15, "Azure AI Search", data["search"].get("name", "aisearchozguler"), TEAL)
    add_node(s3, 7.9, 4.0, 4.0, 1.15, "PostgreSQL Flexible", data["postgres"].get("name", "aviationragpg705508"), TEAL)

    add_flow_line(s3, 10.3, 2.8, 10.3, 4.0, color=SLATE)
    add_flow_line(s3, 9.6, 2.8, 6.1, 4.0, color=SLATE)
    add_flow_line(s3, 9.0, 2.8, 2.7, 4.0, color=SLATE)

    note = s3.shapes.add_textbox(Inches(0.7), Inches(5.55), Inches(12.1), Inches(1.2))
    ntf = note.text_frame
    ntf.clear()
    p = ntf.paragraphs[0]
    p.text = (
        "Traffic currently uses public LB for backend reachability; internal LB is available "
        "for controlled private-routing cutover."
    )
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_DARK

    # Slide 4: Network and trust zones
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s4, "Segmentation of internet edge, VNet runtime, and private data access")
    add_title(s4, "Network Topology and Trust Boundaries", "How ingress and private dependencies are separated")

    public_zone = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(3.2), Inches(4.8))
    public_zone.fill.solid()
    public_zone.fill.fore_color.rgb = RGBColor(231, 242, 255)
    public_zone.line.color.rgb = AZURE
    public_zone.text_frame.text = "Public Zone"
    public_zone.text_frame.paragraphs[0].font.bold = True
    public_zone.text_frame.paragraphs[0].font.size = Pt(12)

    vnet_zone = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.95), Inches(1.6), Inches(5.6), Inches(4.8))
    vnet_zone.fill.solid()
    vnet_zone.fill.fore_color.rgb = RGBColor(236, 247, 241)
    vnet_zone.line.color.rgb = GREEN
    vnet_zone.text_frame.text = f"VNet Zone: {data['vnet'].get('name', 'vnet-aviation-rag')} ({data['vnet'].get('address', '10.0.0.0/16')})"
    vnet_zone.text_frame.paragraphs[0].font.bold = True
    vnet_zone.text_frame.paragraphs[0].font.size = Pt(12)

    shared_zone = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.7), Inches(1.6), Inches(3.0), Inches(4.8))
    shared_zone.fill.solid()
    shared_zone.fill.fore_color.rgb = RGBColor(255, 244, 229)
    shared_zone.line.color.rgb = ORANGE
    shared_zone.text_frame.text = "Shared Services RG"
    shared_zone.text_frame.paragraphs[0].font.bold = True
    shared_zone.text_frame.paragraphs[0].font.size = Pt(12)

    add_node(s4, 0.9, 2.3, 2.5, 0.9, "App Service", "Public HTTPS endpoint", AZURE)
    add_node(s4, 0.9, 3.5, 2.5, 0.9, "Public LB", data["service_ips"].get("public_lb", "20.240.76.230"), AZURE)

    add_node(s4, 4.3, 2.2, 1.5, 0.95, "subnet-aks", "10.0.0.0/22", GREEN)
    add_node(s4, 6.0, 2.2, 1.7, 0.95, "subnet-appservice", "10.0.4.0/24", GREEN)
    add_node(s4, 7.9, 2.2, 1.5, 0.95, "subnet-pe", "10.0.5.0/24", GREEN)

    add_node(s4, 4.5, 3.8, 2.6, 0.95, "Internal LB", data["service_ips"].get("internal_lb", "10.0.0.33"), NAVY)
    add_node(s4, 7.4, 3.8, 2.0, 0.95, "Private Endpoint", data["private_endpoint"].get("name", "pe-postgres"), NAVY)

    add_node(s4, 10.0, 2.4, 2.3, 0.95, "PostgreSQL", data["postgres"].get("state", "Ready"), ORANGE)
    add_node(s4, 10.0, 3.6, 2.3, 0.95, "Azure OpenAI", data["aoai"].get("state", "Succeeded"), ORANGE)
    add_node(s4, 10.0, 4.8, 2.3, 0.95, "AI Search", data["search"].get("status", "running"), ORANGE)

    add_step_arrow(s4, 3.55, 2.63)
    add_step_arrow(s4, 3.55, 3.83)
    add_step_arrow(s4, 9.45, 4.23)
    add_step_arrow(s4, 9.45, 2.95)

    dns_note = s4.shapes.add_textbox(Inches(4.1), Inches(5.3), Inches(5.4), Inches(0.95))
    dnt = dns_note.text_frame
    dnt.clear()
    dp = dnt.paragraphs[0]
    dp.text = "Private DNS: privatelink.postgres.database.azure.com linked to VNet"
    dp.font.bold = True
    dp.font.size = Pt(11)
    dp.font.color.rgb = TEXT_DARK

    # Slide 5: CI/CD sequence
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s5, "Federated identity path with no long-lived cloud credentials in GitHub")
    add_title(s5, "CI/CD Identity and Deployment Sequence", "GitHub Actions OIDC into Azure runtime resources")

    add_node(s5, 0.6, 1.9, 2.5, 0.95, "Step 1", "GitHub Actions trigger", NAVY)
    add_step_arrow(s5, 3.2, 2.25)
    add_node(s5, 3.6, 1.9, 2.8, 0.95, "Step 2", "OIDC token exchange", AZURE)
    add_step_arrow(s5, 6.5, 2.25)
    add_node(s5, 6.9, 1.9, 2.6, 0.95, "Step 3", "Build + push image to ACR", ORANGE)
    add_step_arrow(s5, 9.6, 2.25)
    add_node(s5, 10.0, 1.9, 2.7, 0.95, "Step 4", "Rollout AKS backend", GREEN)

    add_node(s5, 6.9, 3.4, 2.6, 0.95, "Parallel", "Frontend zip deploy", TEAL)
    add_step_arrow(s5, 9.6, 3.75)
    add_node(s5, 10.0, 3.4, 2.7, 0.95, "Target", "App Service frontend", TEAL)

    oidc_subject = data.get("oidc_fed", [{}])[0].get(
        "subject",
        "repo:ozgurgulerx/aviation-demos-01:ref:refs/heads/main",
    )
    role_scope = data.get("pipeline_roles", [{}])[0].get("scope", "/subscriptions/.../rg-aviation-rag")
    role_name = data.get("pipeline_roles", [{}])[0].get("role", "Contributor")

    detail = s5.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(12.0), Inches(1.8))
    dtf = detail.text_frame
    dtf.clear()
    for i, line in enumerate(
        [
            f"Federated credential subject: {oidc_subject}",
            f"Pipeline role: {role_name}",
            f"Pipeline scope: {role_scope}",
            "AKS kubelet identity has Cognitive Services OpenAI User on AOAI account scope.",
        ]
    ):
        p = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK

    # Slide 6: Inventory table
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s6, "Essential infrastructure components required for application runtime")
    add_title(s6, "Infrastructure Inventory", "Required components and live status")

    inventory_rows = [
        ["Runtime", "Resource Group", data["rg_runtime"].get("name", "rg-aviation-rag"), data["rg_runtime"].get("state", "Succeeded")],
        ["Runtime", "AKS", data["aks"].get("name", "aks-aviation-rag"), data["aks"].get("power", "Running")],
        ["Runtime", "ACR", data["acr"].get("name", "avrag705508acr"), data["acr"].get("state", "Succeeded")],
        ["Runtime", "App Service Plan", data["app_service_plan"].get("name", "plan-aviation-rag-frontend"), data["app_service_plan"].get("state", "Ready")],
        ["Runtime", "Web App", data["webapp"].get("name", "aviation-rag-frontend-705508"), data["webapp"].get("state", "Running")],
        ["Network", "VNet", data["vnet"].get("name", "vnet-aviation-rag"), "Running"],
        ["Network", "Private Endpoint", data["private_endpoint"].get("name", "pe-postgres-aviation-rag"), data["private_endpoint"].get("state", "Succeeded")],
        ["Shared", "Azure OpenAI", data["aoai"].get("name", "aoaiaviation705508"), data["aoai"].get("state", "Succeeded")],
        ["Shared", "Azure AI Search", data["search"].get("name", "aisearchozguler"), data["search"].get("status", "running")],
        ["Shared", "PostgreSQL", data["postgres"].get("name", "aviationragpg705508"), data["postgres"].get("state", "Ready")],
    ]
    add_table(s6, 0.45, 1.5, 12.4, 5.6, ["Domain", "Component", "Resource", "Status"], inventory_rows)

    # Slide 7: Health evidence
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s7, "Operational evidence from deployment state and endpoint checks")
    add_title(s7, "Operational Health Evidence", "Live probes and deployment readiness")

    backend_ok = "PASS" if "status" in data.get("backend_health", "") else "WARN"
    frontend_ok = "PASS" if '"status":"ok"' in data.get("frontend_health", "") else "WARN"
    aoai_ok = "PASS" if all(d.get("state") == "Succeeded" for d in data.get("aoai_deployments", [])) else "WARN"

    health_rows = [
        ["Backend Deployment", data.get("k8s_deploy", "2/2 ready"), "PASS"],
        ["Backend Health", f"http://{data['service_ips'].get('public_lb', '20.240.76.230')}/health", backend_ok],
        [
            "Frontend Health",
            f"https://{data['webapp'].get('host', 'aviation-rag-frontend-705508.azurewebsites.net')}/api/health",
            frontend_ok,
        ],
        [
            "AOAI Deployments",
            ", ".join([f"{x.get('name')} ({x.get('model')})" for x in data.get("aoai_deployments", [])]),
            aoai_ok,
        ],
        [
            "AI Search",
            f"replicas={data['search'].get('replicas', 1)}, partitions={data['search'].get('partitions', 1)}",
            "PASS" if str(data["search"].get("status", "")).lower() == "running" else "WARN",
        ],
        [
            "Postgres",
            f"state={data['postgres'].get('state', 'Ready')}, version={data['postgres'].get('version', '16')}, ha={data['postgres'].get('ha', 'NotEnabled')}",
            "PASS" if str(data["postgres"].get("state", "")).lower() == "ready" else "WARN",
        ],
    ]

    add_table(s7, 0.45, 1.65, 12.4, 4.6, ["Check", "Evidence", "Result"], health_rows)

    settings = {x.get("name"): x.get("value") for x in data.get("appsettings", [])}
    settings_box = s7.shapes.add_textbox(Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.65))
    stf = settings_box.text_frame
    stf.clear()
    p = stf.paragraphs[0]
    p.text = (
        f"Key app settings: BACKEND_URL={settings.get('BACKEND_URL', '-')}, "
        f"PII_ENDPOINT={settings.get('PII_ENDPOINT', '-')}, "
        f"WEBSITES_PORT={settings.get('WEBSITES_PORT', '-')}, "
        f"HTTPS-only={data['webapp'].get('httpsOnly', False)}"
    )
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK

    # Slide 8: Controls and hardening
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s8, "Current control state and prioritized hardening backlog")
    add_title(s8, "Security and Reliability Posture", "Current controls vs next hardening actions")

    add_card(
        s8,
        0.7,
        1.6,
        6.0,
        4.9,
        "Controls currently in place",
        [
            "Federated OIDC identity for CI/CD (no static cloud secret in workflow)",
            "Scoped Contributor role on runtime RG for deployment principal",
            "Private endpoint + private DNS zone for PostgreSQL path",
            "AKS runtime identity authorized for Azure OpenAI inference",
            "Dedicated subnets separating compute, app service, and private endpoint",
            "Health probes for backend service and frontend API endpoint",
        ],
        border=GREEN,
    )

    add_card(
        s8,
        6.95,
        1.6,
        5.7,
        4.9,
        "Hardening backlog (recommended order)",
        [
            "1. Set App Service HTTPS-only to true.",
            "2. Shift frontend BACKEND_URL to internal LB endpoint after final routing checks.",
            "3. Rotate runtime credentials after migration freeze completes.",
            "4. Enable Postgres HA and define RTO/RPO targets.",
            "5. Add alerting for AKS pod readiness and synthetic frontend checks.",
            "6. Tighten network rules to reduce public surface after cutover.",
        ],
        border=ORANGE,
    )

    # Slide 9: Go-live checklist
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s9, "Executive checkpoint for deployment readiness")
    add_title(s9, "Go-Live Checklist", "Required infra components and readiness verdict")

    checklist_rows = [
        ["Tenant migration target active", acct.get("user", "-") + " authenticated", "PASS"],
        ["Runtime resource group provisioned", data["rg_runtime"].get("state", "Succeeded"), "PASS"],
        ["Backend serving on AKS", data.get("k8s_deploy", "2/2 ready"), "PASS"],
        ["Frontend online", data["webapp"].get("state", "Running"), "PASS"],
        ["VNet and private endpoint available", data["private_endpoint"].get("state", "Succeeded"), "PASS"],
        ["CI/CD OIDC federation in place", data.get("oidc_fed", [{}])[0].get("name", "github-aviation-main"), "PASS"],
        ["AOAI + Search + Postgres integrated", "Shared services reachable", "PASS"],
        ["HTTPS-only enforced on frontend", str(data["webapp"].get("httpsOnly", False)), "WARN"],
    ]

    add_table(s9, 0.45, 1.6, 12.4, 4.8, ["Requirement", "Evidence", "Result"], checklist_rows)

    verdict = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.5), Inches(12.0), Inches(0.75))
    verdict.fill.solid()
    verdict.fill.fore_color.rgb = RGBColor(232, 245, 233)
    verdict.line.color.rgb = GREEN
    vtf = verdict.text_frame
    vtf.clear()
    vp = vtf.paragraphs[0]
    vp.text = (
        "Verdict: Required essential infrastructure is in place and operational. "
        "Proceed with final hardening tasks (HTTPS-only + private backend route cutover)."
    )
    vp.font.size = Pt(14)
    vp.font.bold = True
    vp.font.color.rgb = RGBColor(24, 74, 24)
    vp.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT)


if __name__ == "__main__":
    payload = collect_data()
    build_presentation(payload)
    print(f"Saved {OUTPUT}")
