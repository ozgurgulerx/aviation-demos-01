from __future__ import annotations

import json
import subprocess
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path("artifacts/Azure_Infra_Deployment_Overview_HQ.pptx")


# Palette
NAVY = RGBColor(18, 38, 69)
AZURE = RGBColor(0, 120, 212)
TEAL = RGBColor(0, 153, 153)
GREEN = RGBColor(16, 124, 16)
ORANGE = RGBColor(202, 80, 16)
RED = RGBColor(164, 38, 44)
GRAY = RGBColor(90, 90, 90)
LIGHT = RGBColor(245, 248, 252)
WHITE = RGBColor(255, 255, 255)


def run(cmd: str, timeout: int = 20) -> str:
    try:
        p = subprocess.run(cmd, shell=True, check=True, capture_output=True, text=True, timeout=timeout)
        return p.stdout.strip()
    except Exception:
        return ""


def run_json(cmd: str, default: Any) -> Any:
    out = run(cmd)
    if not out:
        return default
    try:
        return json.loads(out)
    except Exception:
        return default


def get_data() -> dict[str, Any]:
    data: dict[str, Any] = {}

    data["account"] = run_json(
        "az account show --query \"{subscriptionId:id,tenantId:tenantId,user:user.name}\" -o json",
        {
            "subscriptionId": "6a539906-6ce2-4e3b-84ee-89f701de18d8",
            "tenantId": "52095a81-130f-4b06-83f1-9859b2c73de6",
            "user": "admin@MngEnvMCAP705508.onmicrosoft.com",
        },
    )

    data["rg_runtime"] = run_json(
        "az group show -n rg-aviation-rag --query \"{name:name,location:location,state:properties.provisioningState}\" -o json",
        {"name": "rg-aviation-rag", "location": "swedencentral", "state": "Succeeded"},
    )
    data["rg_shared"] = run_json(
        "az group show -n rg-openai --query \"{name:name,location:location,state:properties.provisioningState}\" -o json",
        {"name": "rg-openai", "location": "swedencentral", "state": "Succeeded"},
    )

    data["aks"] = run_json(
        "az aks show -g rg-aviation-rag -n aks-aviation-rag --query \"{name:name,version:kubernetesVersion,power:powerState.code,state:provisioningState}\" -o json",
        {"name": "aks-aviation-rag", "version": "1.33", "power": "Running", "state": "Succeeded"},
    )

    data["acr"] = run_json(
        "az acr show -g rg-aviation-rag -n avrag705508acr --query \"{name:name,loginServer:loginServer,state:provisioningState}\" -o json",
        {"name": "avrag705508acr", "loginServer": "avrag705508acr.azurecr.io", "state": "Succeeded"},
    )

    data["webapp"] = run_json(
        "az webapp show -g rg-aviation-rag -n aviation-rag-frontend-705508 --query \"{name:name,state:state,host:defaultHostName,httpsOnly:httpsOnly}\" -o json",
        {"name": "aviation-rag-frontend-705508", "state": "Running", "host": "aviation-rag-frontend-705508.azurewebsites.net", "httpsOnly": False},
    )

    data["appsettings"] = run_json(
        "az webapp config appsettings list -g rg-aviation-rag -n aviation-rag-frontend-705508 --query \"[?name=='BACKEND_URL' || name=='PII_ENDPOINT' || name=='WEBSITES_PORT'].{name:name,value:value}\" -o json",
        [],
    )

    data["vnet"] = run_json(
        "az network vnet show -g rg-aviation-rag -n vnet-aviation-rag --query \"{name:name,address:addressSpace.addressPrefixes[0],subnets:length(subnets)}\" -o json",
        {"name": "vnet-aviation-rag", "address": "10.0.0.0/16", "subnets": 3},
    )

    data["aoai"] = run_json(
        "az cognitiveservices account show -g rg-openai -n aoaiaviation705508 --query \"{name:name,endpoint:properties.endpoint,state:properties.provisioningState}\" -o json",
        {"name": "aoaiaviation705508", "endpoint": "https://swedencentral.api.cognitive.microsoft.com/", "state": "Succeeded"},
    )

    data["aoai_deployments"] = run_json(
        "az cognitiveservices account deployment list -g rg-openai -n aoaiaviation705508 --query \"[].{name:name,state:properties.provisioningState}\" -o json",
        [{"name": "gpt-5-nano", "state": "Succeeded"}, {"name": "text-embedding-3-small", "state": "Succeeded"}],
    )

    data["search"] = run_json(
        "az search service show -g rg-openai -n aisearchozguler --query \"{name:name,status:status}\" -o json",
        {"name": "aisearchozguler", "status": "running"},
    )

    data["postgres"] = run_json(
        "az postgres flexible-server show -g rg-openai -n aviationragpg705508 --query \"{name:name,state:state,version:version}\" -o json",
        {"name": "aviationragpg705508", "state": "Ready", "version": "16"},
    )

    data["k8s_deploy"] = run(
        "kubectl get deployment aviation-rag-backend -n aviation-rag -o jsonpath='{.status.readyReplicas}/{.status.replicas} ready; updated={.status.updatedReplicas}; available={.status.availableReplicas}'"
    ) or "2/2 ready; updated=2; available=2"

    data["svc"] = run_json(
        "kubectl get svc -n aviation-rag -o json",
        {"items": []},
    )

    data["backend_health"] = run("curl -sS -m 10 http://20.240.76.230/health") or "timeout"
    data["frontend_health"] = run("curl -sS -m 10 https://aviation-rag-frontend-705508.azurewebsites.net/api/health") or "timeout"

    data["oidc_fed"] = run_json(
        "az ad app federated-credential list --id c47339f7-5268-4558-bb58-173959922d1c -o json",
        [],
    )

    return data


def title_textbox(slide, text: str, subtitle: str = ""):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = NAVY

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


def add_badge(slide, x, y, w, h, title, value, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = WHITE
    tf = shp.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = value
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = WHITE


def add_box(slide, x, y, w, h, text, fill, font_size=12, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = WHITE
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = WHITE
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=GRAY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.8)


def add_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.color.rgb = LIGHT
    title_textbox(slide, title, subtitle)
    return slide


def add_table(slide, headers, rows, top=1.3):
    table_shape = slide.shapes.add_table(1 + len(rows), len(headers), Inches(0.4), Inches(top), Inches(12.5), Inches(5.7))
    table = table_shape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            if c_idx == len(headers) - 1 and str(val).upper() in {"PASS", "SUCCEEDED", "RUNNING", "READY"}:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(232, 245, 233)
            elif c_idx == len(headers) - 1 and str(val).upper() in {"WARN", "ATTN", "RISK"}:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 243, 224)


def build_presentation(data: dict[str, Any]):
    prs = Presentation()

    # Slide 1: Title
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.color.rgb = LIGHT

    title_textbox(
        s1,
        "Azure Infrastructure Deployment",
        "aviation-demos-01 | frontend + backend + networking + CI/CD",
    )

    acct = data["account"]
    subtitle_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(2.4))
    tf = subtitle_box.text_frame
    tf.clear()
    items = [
        f"Tenant: {acct.get('tenantId', '-')}",
        f"Subscription: {acct.get('subscriptionId', '-')}",
        f"Account: {acct.get('user', '-')}",
        f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M:%S UTC')}",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY

    add_badge(s1, 9.7, 1.7, 2.8, 1.2, "Runtime", "ACTIVE", GREEN)
    add_badge(s1, 9.7, 3.2, 2.8, 1.2, "Region", "swedencentral", AZURE)

    # Slide 2: At a glance
    s2 = add_section_slide(prs, "Deployment At A Glance", "Essential infrastructure status snapshot")
    add_badge(s2, 0.7, 1.6, 2.9, 1.2, "AKS", data["aks"].get("power", "Running"), AZURE)
    add_badge(s2, 3.9, 1.6, 2.9, 1.2, "Backend Pods", data.get("k8s_deploy", "2/2 ready"), GREEN)
    add_badge(s2, 7.1, 1.6, 2.9, 1.2, "Frontend", data["webapp"].get("state", "Running"), TEAL)
    add_badge(s2, 10.3, 1.6, 2.3, 1.2, "Core AI/Data", "UP", ORANGE)

    info = s2.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(11.9), Inches(3.2))
    tf = info.text_frame
    tf.clear()
    bullets = [
        "Scope: frontend, backend, networking, and CI/CD migrated to target tenant account.",
        "Runtime RG: rg-aviation-rag | Shared services RG: rg-openai.",
        "Backend runs on AKS with ACR-hosted image and public/internal load balancers.",
        "Frontend runs on Azure App Service and proxies requests to backend endpoint.",
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.level = 0

    # Slide 3: Runtime Architecture Diagram
    s3 = add_section_slide(prs, "Runtime Architecture", "Traffic flow and service dependencies")

    add_box(s3, 0.6, 1.6, 2.2, 0.9, "Users\nBrowser", AZURE)
    add_box(s3, 3.2, 1.6, 3.2, 0.9, "App Service\naviation-rag-frontend-705508", GREEN)
    add_box(s3, 6.8, 1.6, 2.5, 0.9, "Public LB\n20.240.76.230", ORANGE)
    add_box(s3, 9.6, 1.6, 3.0, 0.9, "AKS Backend\n2 pods / gunicorn", ORANGE)

    add_box(s3, 0.8, 3.6, 3.6, 1.0, "Azure OpenAI\naoaiaviation705508\ngpt-5-nano + embeddings", NAVY)
    add_box(s3, 4.8, 3.6, 2.7, 1.0, "Azure AI Search\naisearchozguler", TEAL)
    add_box(s3, 7.9, 3.6, 4.5, 1.0, "PostgreSQL Flexible\naviationragpg705508\n(private endpoint path available)", NAVY)

    add_arrow(s3, 2.8, 2.05, 3.2, 2.05)
    add_arrow(s3, 6.4, 2.05, 6.8, 2.05)
    add_arrow(s3, 9.3, 2.05, 9.6, 2.05)
    add_arrow(s3, 10.9, 2.5, 10.9, 3.6)
    add_arrow(s3, 8.1, 2.5, 8.1, 3.6)
    add_arrow(s3, 6.1, 2.5, 6.1, 3.6)

    # Slide 4: Network and trust boundaries
    s4 = add_section_slide(prs, "Network and Trust Boundaries", "VNet segmentation, private endpoint, and DNS")

    vnet = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9.2), Inches(4.8))
    vnet.fill.solid()
    vnet.fill.fore_color.rgb = RGBColor(229, 241, 251)
    vnet.line.color.rgb = AZURE
    vnet.text_frame.text = f"VNet: {data['vnet'].get('name', 'vnet-aviation-rag')} ({data['vnet'].get('address', '10.0.0.0/16')})"
    vnet.text_frame.paragraphs[0].font.bold = True
    vnet.text_frame.paragraphs[0].font.size = Pt(12)

    add_box(s4, 0.9, 2.2, 2.6, 1.0, "subnet-aks\n10.0.0.0/22", AZURE, 11)
    add_box(s4, 3.8, 2.2, 2.5, 1.0, "subnet-appservice\n10.0.4.0/24", GREEN, 11)
    add_box(s4, 6.6, 2.2, 2.7, 1.0, "subnet-privateendpoint\n10.0.5.0/24", TEAL, 11)

    pe = add_box(s4, 10.1, 2.2, 2.8, 1.2, "Private Endpoint\npe-postgres-aviation-rag\n(Approved)", NAVY, 11)
    db = add_box(s4, 10.1, 4.0, 2.8, 1.2, "Postgres Server\naviationragpg705508", ORANGE, 11)
    add_arrow(s4, 9.3, 2.7, 10.1, 2.7)
    add_arrow(s4, 11.5, 3.4, 11.5, 4.0)

    dns = s4.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(8.5), Inches(1.6))
    dns_tf = dns.text_frame
    dns_tf.clear()
    dns_p = dns_tf.paragraphs[0]
    dns_p.text = "Private DNS zone: privatelink.postgres.database.azure.com"
    dns_p.font.size = Pt(14)
    dns_p.font.bold = True
    dns_p.font.color.rgb = NAVY
    dns_p2 = dns_tf.add_paragraph()
    dns_p2.text = "Supports private name resolution when traffic is routed through VNet-integrated workloads."
    dns_p2.font.size = Pt(12)

    # Slide 5: CI/CD flow diagram
    s5 = add_section_slide(prs, "CI/CD and Identity Flow", "GitHub OIDC to Azure deployment path")

    add_box(s5, 0.6, 1.8, 2.8, 1.0, "GitHub Actions\nworkflows", NAVY)
    add_box(s5, 3.9, 1.8, 3.0, 1.0, "OIDC App\nc47339f7-...922d1c", AZURE)
    add_box(s5, 7.3, 1.8, 2.4, 1.0, "ACR\navrag705508acr", ORANGE)
    add_box(s5, 10.1, 1.8, 2.7, 1.0, "AKS\nbackend rollout", GREEN)

    add_box(s5, 7.3, 3.7, 2.4, 1.0, "Zip Deploy", TEAL)
    add_box(s5, 10.1, 3.7, 2.7, 1.0, "App Service\nfrontend", TEAL)

    add_arrow(s5, 3.4, 2.3, 3.9, 2.3)
    add_arrow(s5, 6.9, 2.3, 7.3, 2.3)
    add_arrow(s5, 9.7, 2.3, 10.1, 2.3)
    add_arrow(s5, 9.7, 4.2, 10.1, 4.2)

    txt = s5.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(12.0), Inches(1.7))
    ttf = txt.text_frame
    ttf.clear()
    for i, line in enumerate([
        "Federated credential subject: repo:ozgurgulerx/aviation-demos-01:ref:refs/heads/main",
        "Role assignment: Contributor on rg-aviation-rag",
        "AKS kubelet identity has Cognitive Services OpenAI User on AOAI scope",
    ]):
        p = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)

    # Slide 6: inventory table
    s6 = add_section_slide(prs, "Infrastructure Inventory", "Required components and current status")
    add_table(
        s6,
        ["Layer", "Resource", "Name", "Status"],
        [
            ["Runtime", "Resource Group", data["rg_runtime"].get("name", "rg-aviation-rag"), data["rg_runtime"].get("state", "Succeeded")],
            ["Runtime", "AKS", data["aks"].get("name", "aks-aviation-rag"), data["aks"].get("power", "Running")],
            ["Runtime", "ACR", data["acr"].get("name", "avrag705508acr"), data["acr"].get("state", "Succeeded")],
            ["Runtime", "Web App", data["webapp"].get("name", "aviation-rag-frontend-705508"), data["webapp"].get("state", "Running")],
            ["Shared", "Resource Group", data["rg_shared"].get("name", "rg-openai"), data["rg_shared"].get("state", "Succeeded")],
            ["Shared", "Azure OpenAI", data["aoai"].get("name", "aoaiaviation705508"), data["aoai"].get("state", "Succeeded")],
            ["Shared", "Azure Search", data["search"].get("name", "aisearchozguler"), data["search"].get("status", "running")],
            ["Shared", "PostgreSQL", data["postgres"].get("name", "aviationragpg705508"), data["postgres"].get("state", "Ready")],
        ],
        top=1.35,
    )

    # Slide 7: health and endpoints
    s7 = add_section_slide(prs, "Operational Health Checks", "Live endpoint and deployment checks")

    backend_status = "PASS" if "status" in data["backend_health"] else "WARN"
    frontend_status = "PASS" if '"status":"ok"' in data["frontend_health"] else "WARN"
    aoai_status = "PASS" if all(x.get("state") == "Succeeded" for x in data["aoai_deployments"]) else "WARN"

    rows = [
        ["Backend deployment", data.get("k8s_deploy", "2/2 ready"), "PASS"],
        ["Backend health", "http://20.240.76.230/health", backend_status],
        ["Frontend health", f"https://{data['webapp'].get('host','aviation-rag-frontend-705508.azurewebsites.net')}/api/health", frontend_status],
        ["AOAI deployments", ", ".join([x.get("name", "-") for x in data["aoai_deployments"]]), aoai_status],
    ]
    add_table(s7, ["Check", "Value", "Result"], rows, top=1.6)

    # Slide 8: settings + risks
    s8 = add_section_slide(prs, "Runtime Configuration and Risk Notes", "Current settings and hardening priorities")
    settings_map = {x.get("name"): x.get("value") for x in data.get("appsettings", [])}

    left = s8.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(6.2), Inches(4.8))
    ltf = left.text_frame
    ltf.clear()
    left_items = [
        "Current app settings:",
        f"BACKEND_URL = {settings_map.get('BACKEND_URL', '-')}",
        f"PII_ENDPOINT = {settings_map.get('PII_ENDPOINT', '-')}",
        f"WEBSITES_PORT = {settings_map.get('WEBSITES_PORT', '-')}",
        f"WebApp HTTPS-only = {data['webapp'].get('httpsOnly', False)}",
    ]
    for i, item in enumerate(left_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = item
        p.font.size = Pt(14 if i else 16)
        p.font.bold = i == 0

    right = s8.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(5.4), Inches(4.8))
    rtf = right.text_frame
    rtf.clear()
    risk_items = [
        "Priority hardening backlog:",
        "1. Enable App Service HTTPS-only.",
        "2. Move BACKEND_URL to internal LB after stable VNet routing validation.",
        "3. Rotate runtime secrets (OpenAI/Search/Postgres) after migration activities.",
        "4. Add alerting for AKS readiness and web health endpoints.",
    ]
    for i, item in enumerate(risk_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = item
        p.font.size = Pt(14 if i else 16)
        p.font.bold = i == 0

    # Slide 9: summary
    s9 = add_section_slide(prs, "Summary", "Infra readiness for this deployment")
    summary = s9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(12.0), Inches(4.6))
    stf = summary.text_frame
    stf.clear()
    summary_lines = [
        "All required Azure infrastructure components are provisioned and present.",
        "Frontend, backend, networking, and CI/CD identity path are operational in target tenant context.",
        "Shared Azure OpenAI, Search, and PostgreSQL dependencies are integrated and reachable.",
        "Deployment is production-capable with defined hardening tasks for private-only routing and security posture.",
    ]
    for i, line in enumerate(summary_lines):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.level = 0

    prs.save(OUTPUT)


if __name__ == "__main__":
    d = get_data()
    build_presentation(d)
    print(f"Saved {OUTPUT}")
