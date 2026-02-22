#!/usr/bin/env python3
"""
Generate a presentation explaining the Fabric Graph Store architecture
for the aviation RAG demo.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/Fabric_Graph_Store_Architecture.pptx")

COLOR_RED = RGBColor(198, 12, 48)
COLOR_DARK = RGBColor(30, 35, 45)
COLOR_LIGHT = RGBColor(245, 247, 250)
COLOR_MID = RGBColor(99, 110, 126)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLUE = RGBColor(33, 100, 180)
COLOR_GREEN = RGBColor(22, 128, 57)
COLOR_AMBER = RGBColor(180, 130, 20)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT
    bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
    top.fill.solid()
    top.fill.fore_color.rgb = COLOR_RED
    top.line.fill.background()


def add_header(slide, title: str, subtitle: str = ""):
    tf = slide.shapes.add_textbox(Inches(0.7), Inches(0.38), Inches(12.0), Inches(0.8)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Segoe UI"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = COLOR_DARK

    if subtitle:
        sf = slide.shapes.add_textbox(Inches(0.72), Inches(1.02), Inches(12.0), Inches(0.45)).text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Segoe UI"
        sr.font.size = Pt(13)
        sr.font.color.rgb = COLOR_MID


def add_footer(slide, text="Aviation RAG - Fabric Graph Store"):
    tf = slide.shapes.add_textbox(Inches(0.7), Inches(7.02), Inches(12.0), Inches(0.3)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(10)
    r.font.color.rgb = COLOR_MID


def add_bullets(slide, bullets, x=0.8, y=1.7, w=12.0, h=5.2, size=18):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = COLOR_DARK


def add_table(slide, headers, rows, top=1.65, height=5.35, font_size=10):
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(0.4), Inches(top), Inches(12.55), Inches(height)
    ).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_RED
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE

    for r_i, row in enumerate(rows, start=1):
        for c_i, value in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = value
            if r_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(235, 239, 245)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLOR_DARK

    return table


def add_node_shape(slide, label: str, x: float, y: float, w: float, h: float, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(
        max(0, color[0] - 40), max(0, color[1] - 40), max(0, color[2] - 40)
    )
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Segoe UI"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = COLOR_WHITE
    return shape


def add_arrow_label(slide, text: str, x: float, y: float, w: float = 1.6, size: int = 9):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = COLOR_RED


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========================================================================
    # Slide 1: Title
    # ========================================================================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s1)
    add_header(s1, "Fabric Graph Store Architecture", "Operational knowledge graph for aviation RAG")
    add_bullets(s1, [
        "Two-layer graph architecture: Operational Graph + Intent Graph",
        "Primary store: Microsoft Fabric KQL/Eventhouse (ops_graph_edges table)",
        "Fallback store: PostgreSQL (same schema, used when Fabric unavailable)",
        "Intent graph drives agentic orchestration: Intent -> Evidence -> Tool routing",
        "Current data: 519 edges, 404 unique nodes, 58 airports, 173 flights, 173 aircraft",
    ], size=17)
    add_footer(s1)

    # ========================================================================
    # Slide 2: Two-Layer Architecture Overview
    # ========================================================================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s2)
    add_header(s2, "Two-Layer Graph Architecture", "Operational graph + Intent graph serve distinct purposes")

    headers = ["Layer", "Purpose", "Storage", "Node Types", "Edge Types"]
    rows = [
        [
            "Operational Graph",
            "Models real-world flight operations: which aircraft fly which routes between which airports",
            "Fabric KQL (primary)\nPostgreSQL (fallback)",
            "Airport, FlightLeg, Tail",
            "DEPARTS, ARRIVES, OPERATES",
        ],
        [
            "Intent Graph",
            "Maps user intents to required evidence and authoritative tools for agentic planning",
            "Fabric Graph endpoint (primary)\nJSON file / built-in default (fallback)",
            "Intent, EvidenceType, Tool",
            "REQUIRES, AUTHORITATIVE_IN, EXPANDS_TO",
        ],
        [
            "Schema-Level (composite)",
            "Defines all possible node/edge types across both layers for routing decisions",
            "schema_provider.py (in-memory)",
            "Airport, Runway, Station, Alternate, Intent, EvidenceType, Tool",
            "REQUIRES, AUTHORITATIVE_IN, EXPANDS_TO, HAS_RUNWAY, HAS_STATION, HAS_ALTERNATE",
        ],
    ]
    table = add_table(s2, headers, rows, font_size=9, top=1.6, height=5.0)
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.6)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.45)
    table.columns[4].width = Inches(2.5)
    add_footer(s2)

    # ========================================================================
    # Slide 3: Operational Graph - Schema & Topology
    # ========================================================================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s3)
    add_header(s3, "Operational Graph: Schema & Data", "ops_graph_edges table (Fabric KQL / PostgreSQL)")

    # Node shapes as visual diagram
    # Airport nodes
    add_node_shape(s3, "Airport\n(58 nodes)", 0.8, 2.0, 2.0, 0.8, COLOR_BLUE)
    # FlightLeg nodes
    add_node_shape(s3, "FlightLeg\n(173 nodes)", 5.5, 2.0, 2.2, 0.8, COLOR_GREEN)
    # Tail nodes
    add_node_shape(s3, "Tail / Aircraft\n(173 nodes)", 10.2, 2.0, 2.3, 0.8, COLOR_AMBER)

    # Edge labels
    add_arrow_label(s3, "--- DEPARTS (173) -->", 2.8, 2.1, 2.7)
    add_arrow_label(s3, "<-- ARRIVES (173) ---", 2.8, 2.5, 2.7)
    add_arrow_label(s3, "<-- OPERATES (173) ---", 7.7, 2.3, 2.5)

    # Data summary table
    headers2 = ["Column", "Type", "Description", "Example Values"]
    rows2 = [
        ["src_type", "string", "Source node type", "Airport, FlightLeg, Tail"],
        ["src_id", "string", "Source node identifier", "IST, LEG0001, N5200TX"],
        ["edge_type", "string", "Relationship type", "DEPARTS, ARRIVES, OPERATES"],
        ["dst_type", "string", "Destination node type", "Airport, FlightLeg, Tail"],
        ["dst_id", "string", "Destination node identifier", "SAW, OSL0014, 4BB868"],
    ]
    table2 = add_table(s3, headers2, rows2, top=3.4, height=2.7, font_size=10)
    table2.columns[0].width = Inches(1.8)
    table2.columns[1].width = Inches(1.2)
    table2.columns[2].width = Inches(4.0)
    table2.columns[3].width = Inches(5.55)

    # Bottom note
    tf = s3.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5)).text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Each flight leg produces 3 edges: Airport --DEPARTS--> FlightLeg --ARRIVES--> Airport, and Tail --OPERATES--> FlightLeg"
    r.font.name = "Segoe UI"
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = COLOR_MID
    add_footer(s3)

    # ========================================================================
    # Slide 4: Network Topology - Hub Analysis
    # ========================================================================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s4)
    add_header(s4, "Operational Graph: Network Topology", "Hub-and-spoke pattern centered on Istanbul airports")

    headers = ["Airport", "Edges", "Role", "Routes"]
    rows = [
        ["IST (Istanbul Airport)", "141", "Primary hub (27% of all edges)", "Domestic Turkey + International (LHR, CDG, DXB, PEK, PVG, BOM, DEL, YYZ, ATH, FCO ...)"],
        ["SAW (Sabiha Gokcen)", "84", "Secondary hub (16% of all edges)", "Domestic Turkey + European (CGN, DTM, HAJ, HAM, STR, TXL, MUC, NUE, DUS, SXF, RIX)"],
        ["CMN (Casablanca)", "4", "International spoke", "IST <-> CMN, SAW <-> CMN"],
        ["CDG (Paris Charles de Gaulle)", "4", "International spoke", "IST <-> CDG, SAW <-> CDG (via LEG routes)"],
        ["LHR (London Heathrow)", "4", "International spoke", "IST <-> LHR"],
        ["DXB (Dubai)", "4", "International spoke", "IST <-> DXB, SAW <-> DXB"],
        ["STR, CGN, HAJ, HAM ...", "2-4 each", "European spoke airports", "SAW-based European routes"],
        ["ADA, ADB, AYT, ERZ ...", "2 each", "Domestic Turkey spokes", "IST-based domestic routes"],
    ]
    table = add_table(s4, headers, rows, font_size=9, top=1.6, height=5.0)
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(0.8)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(6.75)
    add_footer(s4)

    # ========================================================================
    # Slide 5: Intent Graph - Structure
    # ========================================================================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s5)
    add_header(s5, "Intent Graph: Agentic Planning Layer", "Maps user intents to evidence requirements and tool routing")

    # Intent -> Evidence mapping
    headers = ["Intent", "Required Evidence", "Optional Evidence", "Expansion Rule"]
    rows = [
        ["PilotBrief.Departure", "METAR, TAF, NOTAM, RunwayConstraints", "Hazards", "GRAPH: airport -> stations/alternates"],
        ["PilotBrief.Arrival", "METAR, TAF, NOTAM, RunwayConstraints", "-", "GRAPH: airport -> stations/alternates"],
        ["Disruption.Explain", "Hazards, NOTAM", "METAR, SOPClause", "-"],
        ["Policy.Check", "SOPClause", "-", "-"],
        ["Replay.History", "METAR, Hazards", "NOTAM", "-"],
    ]
    table = add_table(s5, headers, rows, font_size=10, top=1.55, height=2.4)
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(4.2)
    table.columns[2].width = Inches(2.4)
    table.columns[3].width = Inches(3.45)

    # Evidence -> Tool mapping
    tf2 = s5.shapes.add_textbox(Inches(0.7), Inches(4.25), Inches(12.0), Inches(0.4)).text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Evidence-to-Tool Authority (AUTHORITATIVE_IN edges)"
    r2.font.name = "Segoe UI"
    r2.font.size = Pt(14)
    r2.font.bold = True
    r2.font.color.rgb = COLOR_DARK

    headers2 = ["Evidence Type", "Primary Tool (Priority 1)", "Secondary Tool (Priority 2)", "Citations Required"]
    rows2 = [
        ["METAR", "KQL (Eventhouse)", "-", "No"],
        ["TAF", "KQL (Eventhouse)", "-", "No"],
        ["NOTAM", "NOSQL", "VECTOR_REG (AI Search)", "No"],
        ["RunwayConstraints", "SQL (Warehouse/PostgreSQL)", "-", "No"],
        ["Hazards", "KQL (Eventhouse)", "-", "No"],
        ["SOPClause", "VECTOR_REG (AI Search)", "-", "Yes"],
    ]
    table2 = add_table(s5, headers2, rows2, top=4.7, height=2.0, font_size=9)
    table2.columns[0].width = Inches(2.2)
    table2.columns[1].width = Inches(3.5)
    table2.columns[2].width = Inches(3.5)
    table2.columns[3].width = Inches(3.35)
    add_footer(s5)

    # ========================================================================
    # Slide 6: Query Execution Flow
    # ========================================================================
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s6)
    add_header(s6, "Graph Query Execution Flow", "How the graph is queried at runtime")

    add_bullets(s6, [
        "1) User query arrives (e.g., 'departure brief for IST')",
        "2) AgenticOrchestrator checks intent graph -> PilotBrief.Departure",
        "3) For PilotBrief intents, GRAPH tool is called first (entity_expansion, hops=2)",
        "4) UnifiedRetriever.query_graph() executes:",
        "     a) Primary: Fabric KQL endpoint -> ops_graph_edges | where src_id in~ ('IST') | take 50",
        "     b) Fallback: PostgreSQL -> SELECT * FROM ops_graph_edges WHERE src_id = 'IST'",
        "5) Expanded entities (connected airports, flight legs, tails) inform subsequent tool calls",
        "6) Dependent calls (KQL for METAR/TAF, SQL for runway data, etc.) execute after graph expansion",
        "7) All evidence is assembled and passed to Azure OpenAI for synthesis",
    ], size=15)
    add_footer(s6)

    # ========================================================================
    # Slide 7: Storage & Fallback Strategy
    # ========================================================================
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s7)
    add_header(s7, "Storage Tiers & Fallback Strategy", "Resilient multi-source graph retrieval")

    headers = ["Component", "Primary Source", "Fallback", "Resolution"]
    rows = [
        [
            "Operational Graph\n(ops_graph_edges)",
            "Fabric KQL/Eventhouse\nvia FABRIC_GRAPH_ENDPOINT\nor FABRIC_KQL_ENDPOINT",
            "PostgreSQL\nops_graph_edges table\n(same 5-column schema)",
            "Automatic: if FABRIC_GRAPH_ENDPOINT\nis not set or returns error,\nfalls back to PostgreSQL query",
        ],
        [
            "Intent Graph\n(intent/evidence/tool mappings)",
            "Fabric Graph endpoint\nPOST { operation: intent_graph_snapshot }",
            "1) JSON file (INTENT_GRAPH_JSON_PATH)\n2) Built-in default (hardcoded)",
            "Priority cascade:\nFabric -> JSON file -> built-in.\nAlways resolves.",
        ],
        [
            "Graph Schema\n(node_types, edge_types)",
            "schema_provider.py\nbuiltin-default",
            "N/A (always available)",
            "Static in-memory definition.\nUsed by query router for\nschema-aware planning.",
        ],
    ]
    table = add_table(s7, headers, rows, font_size=9, top=1.6, height=4.8)
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.3)
    table.columns[2].width = Inches(3.45)
    table.columns[3].width = Inches(3.3)
    add_footer(s7)

    # ========================================================================
    # Slide 8: Key Files & Configuration
    # ========================================================================
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s8)
    add_header(s8, "Key Files & Environment Configuration", "Where graph logic lives in the codebase")

    headers = ["File", "Role", "Key Functions / Classes"]
    rows = [
        ["src/unified_retriever.py", "Core graph query execution", "query_graph(), _query_graph_pg_fallback()"],
        ["src/intent_graph_provider.py", "Intent graph loading (3-tier cascade)", "IntentGraphProvider.load(), IntentGraphSnapshot"],
        ["src/schema_provider.py", "Graph schema definition for routing", "SchemaProvider._graph_schema()"],
        ["src/agentic_orchestrator.py", "Agentic plan with graph expansion step", "AgenticOrchestrator._fallback_plan() (GRAPH tool)"],
        ["src/query_router.py", "Route classification (SQL/SEMANTIC/HYBRID/GRAPH)", "QueryRouter.route()"],
        ["data/.../ops_graph_edges.csv", "Seed data (519 edges, 404 nodes, 58 airports)", "Loaded via scripts/02_load_database.py"],
    ]
    table = add_table(s8, headers, rows, font_size=10, top=1.6, height=3.0)
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(5.05)

    # Env vars
    tf = s8.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(12.0), Inches(0.4)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Environment Variables"
    r.font.name = "Segoe UI"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = COLOR_DARK

    headers2 = ["Variable", "Purpose"]
    rows2 = [
        ["FABRIC_GRAPH_ENDPOINT", "REST or Kusto endpoint for operational graph queries"],
        ["FABRIC_BEARER_TOKEN", "Bearer token for Fabric API authentication (rotatable)"],
        ["FABRIC_KQL_ENDPOINT", "KQL/Eventhouse endpoint (also used for graph when Kusto-style)"],
        ["INTENT_GRAPH_JSON_PATH", "Optional path to static intent graph JSON file"],
    ]
    table2 = add_table(s8, headers2, rows2, top=5.4, height=1.3, font_size=9)
    table2.columns[0].width = Inches(3.5)
    table2.columns[1].width = Inches(9.05)
    add_footer(s8)

    # ========================================================================
    # Slide 9: Current Data Summary
    # ========================================================================
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s9)
    add_header(s9, "Current Graph Data Summary", "Synthetic ops overlay dataset (as of Feb 2026)")

    headers = ["Metric", "Value", "Details"]
    rows = [
        ["Total edges", "519", "3 edges per flight leg (DEPARTS + ARRIVES + OPERATES)"],
        ["Unique nodes", "404", "Across 3 node types"],
        ["Airport nodes", "58", "IST (primary hub, 141 edges), SAW (secondary hub, 84 edges), + 56 spokes"],
        ["FlightLeg nodes", "173", "OSL-prefixed (IST<->SAW shuttles) + LEG-prefixed (domestic + international)"],
        ["Tail/Aircraft nodes", "173", "ICAO 24-bit hex codes (e.g., 4BB868) + N-numbers (e.g., N5200TX)"],
        ["Airports served from IST", "~40", "Domestic Turkey, Europe, Middle East, N. Africa, Central/South Asia, China"],
        ["Airports served from SAW", "~18", "European cities (Germany, Latvia, Baltic) + Middle East"],
        ["Edge types", "3", "DEPARTS (173), ARRIVES (173), OPERATES (173)"],
        ["KQL table", "ops_graph_edges", "520 rows in Fabric Eventhouse (includes one header/meta row)"],
    ]
    table = add_table(s9, headers, rows, font_size=10, top=1.6, height=4.8)
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(8.55)
    add_footer(s9)

    # ========================================================================
    # Save
    # ========================================================================
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
