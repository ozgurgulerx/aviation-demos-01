#!/usr/bin/env python3
"""
Generate a presentation explaining how the intent graph and operational graph
are used during planning and retrieval — deterministic guardrails vs LLM context.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/Graph_Role_In_Agentic_Planning.pptx")

COLOR_RED = RGBColor(198, 12, 48)
COLOR_DARK = RGBColor(30, 35, 45)
COLOR_LIGHT = RGBColor(245, 247, 250)
COLOR_MID = RGBColor(99, 110, 126)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLUE = RGBColor(33, 100, 180)
COLOR_GREEN = RGBColor(22, 128, 57)
COLOR_AMBER = RGBColor(180, 130, 20)
COLOR_TEAL = RGBColor(0, 128, 128)
COLOR_PURPLE = RGBColor(100, 50, 150)


# ── helpers (same style as existing decks) ──────────────────────────────


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT
    bg.line.fill.background()
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
    top.fill.solid()
    top.fill.fore_color.rgb = COLOR_RED
    top.line.fill.background()


def add_header(slide, title: str, subtitle: str = ""):
    tf = slide.shapes.add_textbox(Inches(0.7), Inches(0.38), Inches(12.0), Inches(0.8)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Segoe UI"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = COLOR_DARK
    if subtitle:
        sf = slide.shapes.add_textbox(Inches(0.72), Inches(1.02), Inches(12.0), Inches(0.45)).text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Segoe UI"
        sr.font.size = Pt(13)
        sr.font.color.rgb = COLOR_MID


def add_footer(slide, text="Aviation RAG - Graph Role in Agentic Planning"):
    tf = slide.shapes.add_textbox(Inches(0.7), Inches(7.02), Inches(12.0), Inches(0.3)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(10)
    r.font.color.rgb = COLOR_MID


def add_bullets(slide, bullets, x=0.8, y=1.7, w=12.0, h=5.2, size=18, color=None):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLOR_DARK


def add_table(slide, headers, rows, top=1.65, height=5.35, font_size=10):
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(0.4), Inches(top), Inches(12.55), Inches(height)
    ).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_RED
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE
    for r_i, row in enumerate(rows, start=1):
        for c_i, value in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = value
            if r_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(235, 239, 245)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLOR_DARK
    return table


def add_box(slide, label, x, y, w, h, fill_color, font_size=11, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(max(0, fill_color[0] - 40), max(0, fill_color[1] - 40), max(0, fill_color[2] - 40))
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Segoe UI"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = COLOR_WHITE
    return shape


def add_label(slide, text, x, y, w=2.0, size=10, color=None, bold=False, align=PP_ALIGN.CENTER):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35)).text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or COLOR_DARK


def add_arrow(slide, text, x, y, w=1.8, size=9):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = COLOR_RED


# ── slide builders ──────────────────────────────────────────────────────


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ================================================================
    # Slide 1 — Title / thesis
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "How the Graph Drives Agentic Planning", "Deterministic guardrails + LLM flexibility in a single pipeline")
    add_bullets(s, [
        "The intent graph imposes deterministic evidence-collection paths (hard constraints).",
        "The LLM receives the graph as context and proposes a plan within those constraints.",
        "Post-LLM code patches any gaps the LLM missed, guaranteeing graph completeness.",
        "If the LLM fails entirely, a fully deterministic fallback builds the plan from the graph alone.",
        "The operational graph (ops_graph_edges) is queried at execution time for entity expansion.",
        "Result: the graph is both a deterministic skeleton and an LLM context signal.",
    ], size=17)
    add_footer(s)

    # ================================================================
    # Slide 2 — Four-phase overview table
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Four-Phase Pipeline", "Each phase uses the graph differently")
    headers = ["Phase", "What Happens", "LLM?", "Deterministic?", "Key Code"]
    rows = [
        [
            "1. Deterministic\n    Skeleton",
            "Intent graph lookups:\nIntent -> Evidence -> Tool\n(direct dict lookups, no LLM)",
            "No",
            "Yes",
            "IntentGraphSnapshot\n.required_evidence_for_intent()\n.tools_for_evidence()",
        ],
        [
            "2. LLM Plan\n    Generation",
            "LLM receives intent graph JSON\nas context and proposes plan\n(tool_calls, SQL/KQL, entities)",
            "Yes",
            "No",
            "AgenticOrchestrator\n.create_plan()\npayload[\"intent_graph\"]",
        ],
        [
            "3. Deterministic\n    Post-Patch",
            "Code inspects LLM plan against\ngraph requirements; adds any\nmissing evidence/tool calls",
            "No",
            "Yes",
            "AgenticOrchestrator\n._ensure_required_evidence_calls()",
        ],
        [
            "4. Plan Execution\n    (ops graph)",
            "Operational graph queried for\nentity expansion (airport -> flights,\ntails); results become LLM context",
            "No\n(query)\nYes\n(synthesis)",
            "Yes\n(query)",
            "PlanExecutor._run_call()\nUnifiedRetriever.query_graph()",
        ],
    ]
    table = add_table(s, headers, rows, font_size=9, top=1.55, height=5.0)
    table.columns[0].width = Inches(1.7)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(0.9)
    table.columns[3].width = Inches(1.2)
    table.columns[4].width = Inches(5.25)
    add_footer(s)

    # ================================================================
    # Slide 3 — Phase 1: deterministic skeleton
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Phase 1: Deterministic Skeleton", "Intent graph lookups — no LLM, pure dictionary resolution")

    # Visual: Intent -> Evidence -> Tool boxes
    add_box(s, "PilotBrief\n.Departure", 0.5, 2.0, 1.8, 0.8, COLOR_BLUE, 10)
    add_arrow(s, "REQUIRES -->", 2.35, 2.15, 1.5)
    add_box(s, "METAR", 3.85, 1.7, 1.2, 0.5, COLOR_GREEN, 10)
    add_box(s, "TAF", 3.85, 2.25, 1.2, 0.5, COLOR_GREEN, 10)
    add_box(s, "NOTAM", 3.85, 2.8, 1.2, 0.5, COLOR_GREEN, 10)
    add_box(s, "RunwayConstr.", 3.85, 3.35, 1.2, 0.5, COLOR_GREEN, 10)
    add_box(s, "Hazards (opt)", 3.85, 3.9, 1.2, 0.5, RGBColor(120, 160, 80), 9)
    add_arrow(s, "AUTH_IN -->", 5.1, 2.15, 1.3)
    add_box(s, "KQL", 6.4, 1.7, 1.0, 0.5, COLOR_AMBER, 10)
    add_box(s, "KQL", 6.4, 2.25, 1.0, 0.5, COLOR_AMBER, 10)
    add_box(s, "NOSQL", 6.4, 2.8, 1.0, 0.5, COLOR_AMBER, 10)
    add_box(s, "SQL", 6.4, 3.35, 1.0, 0.5, COLOR_AMBER, 10)
    add_box(s, "KQL", 6.4, 3.9, 1.0, 0.5, COLOR_AMBER, 10)

    add_bullets(s, [
        "required_evidence_for_intent('PilotBrief.Departure')",
        "    -> returns [METAR, TAF, NOTAM, RunwayConstraints] (required) + [Hazards] (optional)",
        "",
        "tools_for_evidence('METAR')  -> ['KQL']   (priority 1)",
        "tools_for_evidence('NOTAM')  -> ['NOSQL', 'VECTOR_REG']   (priority 1, 2)",
        "",
        "This is a pure Python dict lookup. No LLM call. Always the same output for the same intent.",
    ], x=7.8, y=1.7, w=5.2, h=4.0, size=12)

    add_label(s, "Source: intent_graph_provider.py  IntentGraphSnapshot", 0.8, 6.5, w=12.0, size=10, color=COLOR_MID, bold=False, align=PP_ALIGN.LEFT)
    add_footer(s)

    # ================================================================
    # Slide 4 — Phase 2: LLM plan generation
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Phase 2: LLM Plan Generation", "The LLM receives the full intent graph as context — guided, not constrained")

    add_bullets(s, [
        "AgenticOrchestrator.create_plan() sends this payload to Azure OpenAI:",
        "",
        '    payload = {',
        '        "user_query":    "departure brief for IST",',
        '        "intent_graph":  { intents, evidence, tools, requires, authoritative_in, ... },',
        '        "tool_catalog":  { allowed_tools: [GRAPH, KQL, SQL, VECTOR_*, NOSQL] },',
        '        "schemas":       { sql_schema: {...}, kql_schema: {...} },',
        '        "entities":      { airports: ["IST"], ... },',
        '    }',
        "",
        "System prompt instructs the LLM:",
        '    "Use intent_graph as primary guide: Intent -> requires EvidenceType -> authoritative_in Tool."',
        "",
        "The LLM returns a JSON plan: intent classification, tool_calls[], SQL/KQL queries, entity list.",
        "This plan is flexible — the LLM may add extra tool calls or reorder, but must follow the graph.",
    ], size=14, y=1.55)

    add_label(s, "Source: agentic_orchestrator.py  AgenticOrchestrator.create_plan() lines 82-137", 0.8, 6.5, w=12.0, size=10, color=COLOR_MID, bold=False, align=PP_ALIGN.LEFT)
    add_footer(s)

    # ================================================================
    # Slide 5 — Phase 3: deterministic post-patch
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Phase 3: Deterministic Post-Patch", "_ensure_required_evidence_calls() — the graph enforces completeness")

    add_bullets(s, [
        "After the LLM returns its plan, Python code inspects it against the intent graph:",
        "",
        "    for each required evidence type (from intent graph):",
        "        if no tool_call covers this evidence:",
        "            look up authoritative tool from intent graph",
        "            append a new ToolCall to the plan",
        "",
        "Example: LLM returns a plan for PilotBrief.Departure but forgets NOTAM.",
        "    -> _ensure_required_evidence_calls() sees NOTAM is missing",
        "    -> looks up tools_for_evidence('NOTAM') -> ['NOSQL', 'VECTOR_REG']",
        "    -> appends ToolCall(tool='NOSQL', operation='lookup', evidence_type='NOTAM')",
        "",
        "This is a hard override. The LLM cannot skip what the graph mandates.",
        "Coverage items are updated: every required evidence becomes status='planned'.",
    ], size=14, y=1.55)

    add_label(s, "Source: agentic_orchestrator.py  _ensure_required_evidence_calls() lines 279-363", 0.8, 6.5, w=12.0, size=10, color=COLOR_MID, bold=False, align=PP_ALIGN.LEFT)
    add_footer(s)

    # ================================================================
    # Slide 6 — Fallback plan (LLM failure)
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Fallback: Fully Deterministic Plan", "If the LLM call fails, the graph alone builds the entire plan")

    add_bullets(s, [
        "_fallback_plan() executes when the LLM is unavailable or returns invalid JSON.",
        "",
        "Step 1:  _infer_intent(query) — keyword matching, no LLM",
        '            "departure" in query -> PilotBrief.Departure',
        '            "policy" in query    -> Policy.Check',
        "",
        "Step 2:  intent_graph.required_evidence_for_intent(intent_name)",
        "            -> deterministic list of evidence requirements",
        "",
        "Step 3:  For each evidence, intent_graph.tools_for_evidence(name)",
        "            -> deterministic tool assignment (priority-sorted)",
        "",
        "Step 4:  For PilotBrief.* intents, prepend GRAPH entity_expansion call (hops=2)",
        "            -> all subsequent tool calls depend_on the graph call",
        "",
        "Result: a complete AgenticPlan with no LLM involvement whatsoever.",
        "Warning attached: 'LLM routing unavailable; fallback orchestration used.'",
    ], size=14, y=1.50)

    add_label(s, "Source: agentic_orchestrator.py  _fallback_plan() lines 139-230", 0.8, 6.5, w=12.0, size=10, color=COLOR_MID, bold=False, align=PP_ALIGN.LEFT)
    add_footer(s)

    # ================================================================
    # Slide 7 — Phase 4: operational graph at execution time
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Phase 4: Operational Graph at Execution Time", "Entity expansion from ops_graph_edges — then results flow to LLM for synthesis")

    # Visual flow
    add_box(s, "Plan:\nGRAPH call\n(entity_expansion\nhops=2)", 0.5, 2.0, 2.0, 1.2, COLOR_BLUE, 10)
    add_arrow(s, "-- executes -->", 2.55, 2.35, 1.6)
    add_box(s, "query_graph()\nFabric KQL\nor PostgreSQL\nfallback", 4.2, 2.0, 2.0, 1.2, COLOR_GREEN, 10)
    add_arrow(s, "-- returns -->", 6.25, 2.35, 1.5)
    add_box(s, "Expanded\nentities:\nflights, tails,\nconnected airports", 7.8, 2.0, 2.0, 1.2, COLOR_AMBER, 10)
    add_arrow(s, "-- unblocks -->", 9.85, 2.35, 1.5)
    add_box(s, "Dependent\ncalls:\nKQL, SQL,\nVECTOR, NOSQL", 11.4, 2.0, 1.7, 1.2, COLOR_PURPLE, 10)

    add_bullets(s, [
        "PlanExecutor respects depends_on: GRAPH call runs first, other calls wait.",
        "",
        "Fabric KQL path:  ops_graph_edges | where src_id in~ ('IST') or dst_id in~ ('IST') | take 50",
        "PostgreSQL path:  SELECT * FROM ops_graph_edges WHERE UPPER(src_id) IN ('IST') LIMIT 50",
        "",
        "Graph results (airport -> flight legs -> tails) become part of source_results['GRAPH'].",
        "All source results (GRAPH + KQL + SQL + VECTOR + NOSQL) are composed into context_text.",
        "Context text is sent to Azure OpenAI for final answer synthesis.",
        "",
        "The LLM sees graph data as retrieval results for synthesis — not for planning.",
    ], x=0.8, y=3.6, w=12.0, h=3.5, size=13)

    add_label(s, "Source: plan_executor.py lines 289-292, unified_retriever.py query_graph() lines 1202-1243", 0.8, 6.5, w=12.0, size=10, color=COLOR_MID, bold=False, align=PP_ALIGN.LEFT)
    add_footer(s)

    # ================================================================
    # Slide 8 — Summary: who controls what
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Summary: Who Controls What", "The graph sets the floor; the LLM adds intelligence above it")

    headers = ["Aspect", "Controlled By", "Mechanism", "Can the LLM Override?"]
    rows = [
        [
            "Which evidence types\nare required",
            "Intent Graph\n(deterministic)",
            "REQUIRES edges:\nIntent -> EvidenceType",
            "No. Post-patch adds any\nmissing calls back.",
        ],
        [
            "Which tool serves\neach evidence",
            "Intent Graph\n(deterministic)",
            "AUTHORITATIVE_IN edges:\nEvidence -> Tool (priority sorted)",
            "No. Post-patch enforces\nprimary tool assignment.",
        ],
        [
            "Whether GRAPH expansion\nruns first",
            "Intent Graph\n(expansion_rules)",
            "PilotBrief.* intents always\nget GRAPH call_1 prepended",
            "No. Hardcoded in\n_fallback_plan().",
        ],
        [
            "Intent classification",
            "LLM (primary)\nKeyword heuristic (fallback)",
            "LLM returns intent.name\nFallback: _infer_intent()",
            "LLM proposes;\nfallback is deterministic.",
        ],
        [
            "SQL/KQL query text",
            "LLM (via SQLWriter /\nKQLWriter)",
            "LLM generates query\nfrom schema + entities",
            "Yes. LLM has full\nquery authorship.",
        ],
        [
            "Entity extraction\n(airports, flights)",
            "LLM (primary)\nRegex (fallback)",
            "LLM extracts from query;\ncode extracts ICAO codes",
            "LLM proposes;\ncode merges both.",
        ],
        [
            "Answer synthesis",
            "LLM",
            "Azure OpenAI receives all\nretrieval results as context",
            "Yes. LLM has full\nauthorship of final answer.",
        ],
    ]
    table = add_table(s, headers, rows, font_size=9, top=1.55, height=5.2)
    table.columns[0].width = Inches(2.3)
    table.columns[1].width = Inches(2.3)
    table.columns[2].width = Inches(4.0)
    table.columns[3].width = Inches(3.95)
    add_footer(s)

    # ================================================================
    # Slide 9 — Design rationale
    # ================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_header(s, "Why This Design", "Deterministic guardrails + LLM flexibility = reliable aviation intelligence")

    add_bullets(s, [
        "Safety: Aviation briefings must include all required evidence (METAR, TAF, NOTAM, ...).",
        "    The graph guarantees completeness — even if the LLM hallucinates or fails.",
        "",
        "Flexibility: The LLM can add extra tool calls, reorder priorities, generate custom SQL/KQL.",
        "    This allows the system to handle novel queries the graph alone cannot anticipate.",
        "",
        "Resilience: Three-tier fallback (Fabric -> JSON -> built-in) for the intent graph,",
        "    and two-tier fallback (Fabric KQL -> PostgreSQL) for the operational graph.",
        "    The system never fails to produce a plan.",
        "",
        "Auditability: Every plan records graph_source, coverage, warnings, and source traces.",
        "    Operators can see exactly which evidence was planned, which was post-patched,",
        "    and which was missing — with full datastore attribution.",
        "",
        "Evolvability: New intents, evidence types, or tool mappings are added to the intent graph",
        "    (Fabric or JSON) — no code changes needed for the planning pipeline.",
    ], size=14, y=1.55)
    add_footer(s)

    # ── save ────────────────────────────────────────────────────────
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
