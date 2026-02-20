from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


OUTPUT_PATH = "docs/Azure_Infra_Deployment_Overview.pptx"


def set_title(run):
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(22, 49, 82)


def set_subtitle(run):
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(64, 64, 64)


def add_title_and_bullets(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
    return slide


def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    table_shape = slide.shapes.add_table(
        rows=1 + len(rows),
        cols=len(headers),
        left=Inches(0.4),
        top=Inches(1.4),
        width=Inches(12.5),
        height=Inches(5.6),
    )
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)

    return slide


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Runtime Architecture"

    def box(x, y, w, h, text, color):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.color.rgb = RGBColor(255, 255, 255)
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        return shp

    c_blue = RGBColor(0, 120, 212)
    c_green = RGBColor(16, 124, 16)
    c_orange = RGBColor(202, 80, 16)
    c_purple = RGBColor(111, 66, 193)
    c_teal = RGBColor(0, 153, 153)

    box(0.4, 1.4, 2.6, 0.9, "Users / Browser", c_blue)
    box(3.4, 1.4, 3.1, 0.9, "App Service\naviation-rag-frontend-705508", c_green)
    box(7.0, 1.4, 2.7, 0.9, "AKS Service LB\n20.240.76.230", c_orange)
    box(10.1, 1.4, 2.4, 0.9, "AKS Backend Pods\n2 replicas", c_orange)

    box(0.8, 3.2, 3.3, 0.9, "Azure OpenAI\naoaiaviation705508\n(gpt-5-nano, embedding)", c_purple)
    box(4.6, 3.2, 2.8, 0.9, "Azure AI Search\naisearchozguler", c_teal)
    box(7.9, 3.2, 4.2, 0.9, "PostgreSQL Flexible Server\naviationragpg705508\n(via Private Endpoint)", c_purple)

    # connectors (simple lines)
    connectors = [
        (3.0, 1.85, 3.4, 1.85),
        (6.5, 1.85, 7.0, 1.85),
        (9.7, 1.85, 10.1, 1.85),
        (10.9, 2.3, 10.9, 3.2),
        (9.0, 2.3, 9.0, 3.2),
        (6.0, 2.3, 6.0, 3.2),
    ]
    for x1, y1, x2, y2 in connectors:
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = RGBColor(80, 80, 80)
        line.line.width = Pt(1.5)


if __name__ == "__main__":
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_tf = slide.shapes.title.text_frame
    title_tf.clear()
    title_run = title_tf.paragraphs[0].add_run()
    title_run.text = "Azure Infrastructure Overview"
    set_title(title_run)

    subtitle = slide.placeholders[1].text_frame
    subtitle.clear()
    sub_run = subtitle.paragraphs[0].add_run()
    sub_run.text = (
        "Deployment: aviation-demos-01\\n"
        "Tenant: MngEnvMCAP705508.onmicrosoft.com\\n"
        f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')} UTC"
    )
    set_subtitle(sub_run)

    add_title_and_bullets(
        prs,
        "Scope and Essentials",
        [
            "Scope includes frontend, backend, networking, CI/CD, and essential runtime infrastructure.",
            "Data stack migration handled separately; runtime is integrated to shared AI/Search/Postgres services.",
            "Primary runtime resource group: rg-aviation-rag (swedencentral).",
            "Shared services resource group: rg-openai (swedencentral).",
        ],
    )

    add_architecture_slide(prs)

    add_table_slide(
        prs,
        "Resource Groups and Core Services",
        ["Layer", "Resource", "Name", "Status"],
        [
            ["Runtime", "Resource Group", "rg-aviation-rag", "Succeeded"],
            ["Shared", "Resource Group", "rg-openai", "Succeeded"],
            ["Compute", "AKS", "aks-aviation-rag", "Running / Succeeded"],
            ["Container", "ACR", "avrag705508acr", "Succeeded"],
            ["Web", "App Service", "aviation-rag-frontend-705508", "Running"],
        ],
    )

    add_table_slide(
        prs,
        "Networking Topology",
        ["Component", "Name/Value", "Purpose", "Status"],
        [
            ["VNet", "vnet-aviation-rag (10.0.0.0/16)", "Network boundary", "Succeeded"],
            ["AKS Subnet", "subnet-aks (10.0.0.0/22)", "AKS nodes/pods", "Active"],
            ["AppService Subnet", "subnet-appservice (10.0.4.0/24)", "Web app integration", "Active"],
            ["PE Subnet", "subnet-privateendpoint (10.0.5.0/24)", "Private endpoints", "Active"],
            ["Private Endpoint", "pe-postgres-aviation-rag", "Private path to Postgres", "Approved / Succeeded"],
            ["Private DNS", "privatelink.postgres.database.azure.com", "Name resolution for PE", "Configured"],
        ],
    )

    add_table_slide(
        prs,
        "AKS Runtime Services",
        ["Service", "Type", "Endpoint", "Role"],
        [
            ["aviation-rag-backend", "ClusterIP", "10.1.153.149:5001", "In-cluster backend service"],
            ["aviation-rag-backend-internal", "LoadBalancer (internal)", "10.0.0.33:80", "Private LB access"],
            ["aviation-rag-backend-lb", "LoadBalancer (public)", "20.240.76.230:80", "Public LB fallback/ingress"],
            ["Backend Deployment", "AKS Deployment", "2/2 ready replicas", "API runtime availability"],
        ],
    )

    add_table_slide(
        prs,
        "Data and AI Dependencies",
        ["Service", "Name", "Key Config", "Status"],
        [
            ["Azure OpenAI", "aoaiaviation705508", "Endpoint: swedencentral.api.cognitive.microsoft.com", "Succeeded"],
            ["AOAI Deployment", "gpt-5-nano", "Model version: 2025-08-07", "Succeeded"],
            ["AOAI Deployment", "text-embedding-3-small", "Embedding model", "Succeeded"],
            ["Azure Search", "aisearchozguler", "1 replica / 1 partition", "running"],
            ["PostgreSQL Flexible", "aviationragpg705508", "Version 16, state Ready", "Ready"],
        ],
    )

    add_table_slide(
        prs,
        "CI/CD and Identity",
        ["Area", "Implementation", "Value", "Status"],
        [
            ["OIDC App", "Service Principal", "c47339f7-5268-4558-bb58-173959922d1c", "Present"],
            ["Federated Credential", "GitHub Actions OIDC", "repo:ozgurgulerx/aviation-demos-01:ref:refs/heads/main", "Configured"],
            ["RG Role", "Contributor", "Assigned on rg-aviation-rag", "Assigned"],
            ["AKS->AOAI Role", "Cognitive Services OpenAI User", "Assigned to AKS kubelet identity", "Assigned"],
            ["Workflows", "GitHub Actions", "deploy-backend / deploy-frontend / infra-health-check / migrate-database", "Present"],
        ],
    )

    add_title_and_bullets(
        prs,
        "Operational Validation",
        [
            "Frontend health endpoint: /api/health returns status=ok.",
            "Backend health endpoint: /health returns status=ok.",
            "Backend pods currently stable at 2 replicas in AKS.",
            "Current frontend BACKEND_URL points to public backend LB for runtime stability.",
            "Recommended hardening: enable App Service HTTPS-only and complete private-only backend routing.",
        ],
    )

    add_title_and_bullets(
        prs,
        "Summary",
        [
            "Essential Azure infrastructure for this deployment is provisioned and running.",
            "Frontend, backend, networking, and CI/CD are in place under the target tenant/account.",
            "Shared AI/search/database services are connected and operational.",
        ],
    )

    prs.save(OUTPUT_PATH)
    print(f"Saved {OUTPUT_PATH}")
