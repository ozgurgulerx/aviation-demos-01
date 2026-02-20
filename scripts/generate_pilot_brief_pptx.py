#!/usr/bin/env python3
"""Generate the Aviation Demo - Context Engineering PPTX."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/Aviation_Demo_Context_Engineering_Retrieval_Detailed.pptx")

COLOR_RED = RGBColor(198, 12, 48)
COLOR_DARK = RGBColor(30, 35, 45)
COLOR_LIGHT = RGBColor(245, 247, 250)
COLOR_MID = RGBColor(86, 98, 116)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_PANEL = RGBColor(233, 238, 247)


def add_bg(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT
    bg.line.fill.background()

    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.26)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = COLOR_RED
    top.line.fill.background()


def add_header(slide, title: str, subtitle: str = ""):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.42), Inches(12.0), Inches(0.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.size = Pt(31)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.67), Inches(1.08), Inches(12.1), Inches(0.45))
        st = sub.text_frame
        st.clear()
        sp = st.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Segoe UI"
        sp.font.size = Pt(14)
        sp.font.color.rgb = COLOR_MID


def add_footer(slide, text: str = "Aviation Demo - Context Engineering"):
    f = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.28))
    tf = f.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_MID


def add_bullets(slide, x, y, w, h, bullets, font_size=17, line_spacing=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = "Segoe UI"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(7)
        p.line_spacing = line_spacing


def add_table(slide, headers, rows, col_widths, top=1.62, height=5.45, font_size=10):
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(0.38), Inches(top), Inches(12.58), Inches(height)
    ).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_RED
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(237, 241, 248)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLOR_DARK

    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Inches(width)


def add_store_detail_slide(slide, store_name: str, holds: list[str], how: list[str], why: list[str]):
    panel1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(4.0), Inches(4.9))
    panel2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.83), Inches(1.7), Inches(4.0), Inches(4.9))
    panel3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.06), Inches(1.7), Inches(3.65), Inches(4.9))
    for panel in (panel1, panel2, panel3):
        panel.fill.solid()
        panel.fill.fore_color.rgb = COLOR_PANEL
        panel.line.color.rgb = COLOR_RED
        panel.line.width = Pt(1.1)

    t1 = panel1.text_frame
    t1.clear()
    t1.paragraphs[0].text = "What It Holds"
    t1.paragraphs[0].font.name = "Segoe UI"
    t1.paragraphs[0].font.bold = True
    t1.paragraphs[0].font.size = Pt(14)
    t1.paragraphs[0].font.color.rgb = COLOR_DARK
    for item in holds:
        p = t1.add_paragraph()
        p.text = item
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK

    t2 = panel2.text_frame
    t2.clear()
    t2.paragraphs[0].text = "How It Is Used"
    t2.paragraphs[0].font.name = "Segoe UI"
    t2.paragraphs[0].font.bold = True
    t2.paragraphs[0].font.size = Pt(14)
    t2.paragraphs[0].font.color.rgb = COLOR_DARK
    for item in how:
        p = t2.add_paragraph()
        p.text = item
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK

    t3 = panel3.text_frame
    t3.clear()
    t3.paragraphs[0].text = "Why Important for Airline"
    t3.paragraphs[0].font.name = "Segoe UI"
    t3.paragraphs[0].font.bold = True
    t3.paragraphs[0].font.size = Pt(14)
    t3.paragraphs[0].font.color.rgb = COLOR_DARK
    for item in why:
        p = t3.add_paragraph()
        p.text = item
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK

    add_header(slide, store_name, "Detailed datastore role in retrieval pipeline")
    add_footer(slide)


def add_data_shape_slide(
    slide,
    title: str,
    format_text: str,
    fields: list[str],
    sample: list[str],
    note: str = "",
):
    add_bg(slide)
    add_header(slide, f"Data Shape - {title}", "How this source is structured in the demo")

    top = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(1.65), Inches(12.1), Inches(0.8))
    top.fill.solid()
    top.fill.fore_color.rgb = COLOR_PANEL
    top.line.color.rgb = COLOR_RED
    tf = top.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Format: {format_text}"
    p.font.name = "Segoe UI"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_DARK

    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(2.62), Inches(5.9), Inches(3.95))
    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.82), Inches(2.62), Inches(5.9), Inches(3.95))
    for panel in (left, right):
        panel.fill.solid()
        panel.fill.fore_color.rgb = COLOR_PANEL
        panel.line.color.rgb = COLOR_RED
        panel.line.width = Pt(1.1)

    ltf = left.text_frame
    ltf.clear()
    p0 = ltf.paragraphs[0]
    p0.text = "Key Fields"
    p0.font.name = "Segoe UI"
    p0.font.bold = True
    p0.font.size = Pt(14)
    p0.font.color.rgb = COLOR_DARK
    for item in fields:
        p = ltf.add_paragraph()
        p.text = item
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK

    rtf = right.text_frame
    rtf.clear()
    p0 = rtf.paragraphs[0]
    p0.text = "Example Record"
    p0.font.name = "Segoe UI"
    p0.font.bold = True
    p0.font.size = Pt(14)
    p0.font.color.rgb = COLOR_DARK
    for line in sample:
        p = rtf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_DARK

    if note:
        n = slide.shapes.add_textbox(Inches(0.67), Inches(6.7), Inches(12.0), Inches(0.3)).text_frame
        n.clear()
        np = n.paragraphs[0]
        np.text = note
        np.font.name = "Segoe UI"
        np.font.size = Pt(10)
        np.font.color.rgb = COLOR_MID

    add_footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s1)
    add_header(
        s1,
        "Aviation Demo - Context Engineering",
        "Retrieval-centric architecture with agentic routing across Microsoft Fabric and Azure",
    )
    add_bullets(
        s1,
        0.75,
        2.0,
        12.0,
        3.6,
        [
            "Goal: produce high-confidence pilot brief answers with source-attributed evidence.",
            "Design: one query can blend telemetry, relational facts, graph paths, and unstructured narratives.",
            "Scope: Agent Framework runtime + Azure AI Foundry IQ orchestration + Microsoft Fabric data plane.",
            "Constraint policy: deterministic retrieval first, reasoning second, explicit confidence in final output.",
        ],
        font_size=20,
    )
    add_footer(s1)

    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s2)
    add_header(s2, "Why This Multi-Store Structure", "Each store is optimized for a different retrieval function")
    add_bullets(
        s2,
        0.78,
        1.8,
        12.0,
        4.9,
        [
            "Microsoft Fabric Eventhouse (KQL): live telemetry + recent event context.",
            "Microsoft Fabric Warehouse: stable relational joins, KPIs, deterministic reporting.",
            "Microsoft Fabric Lakehouse: raw landing + batch curation.",
            "Azure AI Search: vector/hybrid retrieval for narratives and regulatory documents.",
            "Microsoft Fabric Graph (preview): relationship routing and multi-hop context expansion.",
            "Reasoning: this split minimizes latency/quality tradeoffs while keeping retrieval auditable.",
        ],
        font_size=17,
    )
    add_footer(s2)

    # Slide 3
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s3)
    add_header(s3, "Acronym Glossary", "Long forms used in architecture and query examples")
    acronym_rows = [
        ["RAG", "Retrieval-Augmented Generation", "Grounded answers with explicit evidence."],
        ["KQL", "Kusto Query Language", "Low-latency event/time-window analytics in Eventhouse."],
        ["SQL", "Structured Query Language", "Deterministic relational retrieval in Warehouse."],
        ["NOTAM", "Notice to Air Missions", "Operational notices that impact dispatch/briefing."],
        ["METAR", "Meteorological Aerodrome Report", "Current airport weather observations."],
        ["TAF", "Terminal Aerodrome Forecast", "Airport weather forecast for briefing horizon."],
        ["SIGMET", "Significant Meteorological Information", "Hazard advisories affecting route safety."],
        ["AIRMET", "Airmen’s Meteorological Information", "Lower-severity weather advisories for operations."],
        ["PIREP", "Pilot Report", "Pilot-observed weather and turbulence signals."],
        ["ASRS", "Aviation Safety Reporting System", "Narrative corpus for similar-case retrieval."],
        ["NTSB", "National Transportation Safety Board", "Incident/accident records for lessons learned."],
        ["AD", "Airworthiness Directive", "Mandatory continuing airworthiness requirements."],
        ["MEL", "Minimum Equipment List", "Dispatch rules for inoperative equipment."],
    ]
    add_table(s3, ["Acronym", "Long Form", "Why It Matters"], acronym_rows, [1.3, 4.6, 6.68], top=1.54, height=5.8, font_size=9.2)
    add_footer(s3)

    # Slide 4
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s4)
    add_header(s4, "Data Inventory and Operational Purpose", "What data exists and why each source is aviation-critical")
    data_rows = [
        ["AviationWeather", "METAR/TAF/SIGMET/AIRMET/PIREP", "Flight risk and pre-departure safety context", "Microsoft Fabric Eventhouse + Microsoft Fabric Lakehouse"],
        ["OpenSky", "State vectors, tracks", "Live traffic state and nearby traffic pressure", "Microsoft Fabric Eventhouse + Microsoft Fabric Lakehouse"],
        ["OurAirports", "Airports/runways/navaids/frequencies", "Airport topology and runway suitability checks", "Microsoft Fabric Warehouse + Microsoft Fabric Graph (preview)"],
        ["OpenFlights", "Routes/airlines/airports", "Network context and alternate routing options", "Microsoft Fabric Warehouse + Microsoft Fabric Graph (preview)"],
        ["ASRS", "Safety narratives", "Case-based reasoning for similar disruptions", "Azure AI Search (idx_ops_narratives)"],
        ["NTSB", "Incident narratives + metadata", "Historical risk analogs and mitigation evidence", "Azure AI Search (idx_ops_narratives) + Microsoft Fabric Lakehouse"],
        ["EASA AD", "Regulatory notices + PDFs", "Fleet compliance and dispatch restrictions", "Azure AI Search (idx_regulatory) + Microsoft Fabric Lakehouse"],
        ["Synthetic Ops Overlay", "Crew/gate/baggage/turnaround/MEL", "Operational realism for multi-agent decisions", "Microsoft Fabric Warehouse + Microsoft Fabric Eventhouse + Microsoft Fabric Graph (preview)"],
    ]
    add_table(
        s4,
        ["Source", "Data Type", "Why Important for Airline", "Primary Datastore"],
        data_rows,
        [1.8, 3.0, 3.9, 3.88],
        top=1.55,
        height=5.8,
        font_size=9.2,
    )
    add_footer(s4)

    # Slide 5-11: data shape slides
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_data_shape_slide(
        s5,
        "AviationWeather (METAR/TAF/SIGMET/AIRMET/PIREP)",
        "CSV.gz and XML.gz caches; near real-time snapshots.",
        [
            "raw_text, station_id, observation_time, latitude, longitude",
            "temp_c, dewpoint_c, wind_dir_degrees, wind_speed_kt, visibility_statute_mi",
            "hazard feeds add valid_time_from/to, severity, hazard, geometry points",
        ],
        [
            "station_id=UNBB; observation_time=2026-02-18T15:37:00Z",
            "temp_c=3; wind_dir_degrees=230; wind_speed_kt=19",
            "hazard=CONVECTIVE; severity=SEV; type=SIGMET",
            "polygon=-85.012:43.145;-82.894:42.953;...",
        ],
    )

    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_data_shape_slide(
        s6,
        "OpenSky Live States and Flights",
        "REST JSON objects and arrays.",
        [
            "states payload: time + states[list]",
            "state vector: icao24, callsign, last_contact, lon, lat, baro_altitude, velocity",
            "flight payload: estDepartureAirport, estArrivalAirport, firstSeen, lastSeen",
            "track payload: startTime, endTime, path[list of points]",
        ],
        [
            "time=1771519645",
            "states[0]=['39de4f','TVF93JE',...,196.72,25.89,...]",
            "callsign='THY6047'; estDepartureAirport='LTFM'",
            "path[0]=[1771517590,41.9796,22.2923,5486,90,false]",
        ],
    )

    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_data_shape_slide(
        s7,
        "OurAirports and OpenFlights Network",
        "CSV and .dat (CSV-like quoted rows).",
        [
            "OurAirports airports: ident, type, name, lat/lon, iso_country, gps_code",
            "OurAirports runways: airport_ident, length_ft, surface, le/he heading",
            "OpenFlights routes: airline, source_airport, dest_airport, stops, equipment",
            "OpenFlights airports: airport_id, city, country, iata, icao, timezone",
        ],
        [
            "OurAirports: ident=00A; type=heliport; municipality=Bensalem",
            "Runway: airport_ident=00A; length_ft=80; surface=ASPH-G",
            "Route row: airline='2B', src='AER', dst='KZN', stops='0'",
            "Airport row: id=1, name='Goroka Airport', iata='GKA'",
        ],
    )

    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_data_shape_slide(
        s8,
        "ASRS Narratives",
        "Processed JSONL and SQL table records.",
        [
            "asrs_report_id, event_date, location, aircraft_type, flight_phase",
            "narrative_type, title, report_text, source, raw_json",
            "Long free-text narrative body used for vector chunking",
        ],
        [
            "asrs_report_id='100007'; event_date='1988-12-01'",
            "location='LIZ, CA'; flight_phase='Initial Climb'",
            "title='ASRS | 1988-12-01 | ...'",
            "report_text='MY WIFE AND I WERE DEPARTING ...'",
        ],
    )

    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_data_shape_slide(
        s9,
        "NTSB Accident Archive",
        "ZIP containing Microsoft Access MDB databases (raw archive).",
        [
            "Current repo stores raw: c2-avall.zip, c2-avall.mdb, c2-PRE1982.zip",
            "Table-based relational records inside MDB (not CSV directly)",
            "Canonical parse target: event_id, event_date, location, aircraft metadata",
            "Canonical parse target: injury/fatality indicators, probable cause, narrative",
        ],
        [
            "raw_file='c2-avall.mdb'",
            "container='Access MDB' (table-oriented format)",
            "parse_target.event_date='YYYY-MM-DD'",
            "parse_target.narrative='free text probable cause/summary'",
        ],
        note="Note: NTSB is currently staged as raw MDB in this workspace; canonical extraction is a planned parsing step.",
    )

    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_data_shape_slide(
        s10,
        "EASA AD and NOTAM Corpus",
        "CSV metadata + PDF documents + JSONL NOTAM records.",
        [
            "EASA metadata: class_number, issue_date, effective_date, subject, pdf_url",
            "NOTAM JSONL: facilityDesignator, notamNumber, issue/start/end dates, icaoMessage",
            "PDF body text extracted/chunked for regulatory retrieval",
        ],
        [
            "class_number='CF-2026-08'; effective_date='2026-02-19'",
            "subject='Overhead Stowage Compartment ...'",
            "facilityDesignator='LTBA'; notamNumber='G1555/12'",
            "icaoMessage='Q) LTBB/... A) LTBA B) ... C) PERM ...'",
        ],
    )

    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_data_shape_slide(
        s11,
        "Synthetic Operational Overlay and Schedule",
        "CSV synthetic tables + schedule ZIP/CSV extracts.",
        [
            "ops_flight_legs: flight_no, origin_iata, dest_iata, dep/arr times, tailnum",
            "ops_turnaround_milestones: milestone, event_ts_utc, status, delay_cause_code",
            "ops_crew_rosters: duty windows, role, cumulative_duty_hours, legality_risk_flag",
            "ops_graph_edges: src_type, src_id, edge_type, dst_type, dst_id",
            "schedule_delay_causes: carrier_ct, weather_ct, nas_ct, late_aircraft_delay",
        ],
        [
            "leg_id='OSL0014'; flight_no='PGT161K'; origin='SAW'; dest='IST'",
            "milestone='GATE_OPEN'; status='done'; delay_cause_code='NONE'",
            "duty_id='DUTY-CAP002-0001'; legality_risk_flag='0'",
            "edge: Airport SAW -DEPARTS-> FlightLeg OSL0014",
        ],
    )

    # Slide 12
    s12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s12)
    add_header(s12, "Context-to-Datastore Mapping", "Chosen store, retrieval mode, and selection rationale")
    mapping_rows = [
        ["Live weather + hazards", "Microsoft Fabric Eventhouse (KQL)", "Sub-minute event windows and anomaly scans", "Fast rolling brief updates need event-native queries"],
        ["Ops relational state (crew, gate, MEL)", "Microsoft Fabric Warehouse", "Deterministic SQL joins with constraints", "Dispatch decisions require auditable KPI math"],
        ["Raw evidence and replay", "Microsoft Fabric Lakehouse", "Immutable landing zone + Spark transforms", "Traceability and reprocessing for model drift/debugging"],
        ["Narratives and regulations", "Azure AI Search (vector + hybrid)", "Semantic + lexical ranking with metadata filters", "Safety/regulatory answers are document-heavy"],
        ["Dependency/context graph", "Microsoft Fabric Graph (preview)", "Multi-hop traversal for route/asset dependencies", "Complex disruptions require relationship reasoning"],
    ]
    add_table(
        s12,
        ["Context Family", "Chosen Datastore", "Retrieval Mode", "Reasoning for Choice"],
        mapping_rows,
        [2.15, 3.3, 3.25, 3.88],
        top=1.68,
        height=5.35,
        font_size=9.6,
    )
    add_footer(s12)

    # Slide 13
    s13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s13)
    add_header(s13, "Complex Query Examples for Agentic Retrieval", "Each scenario blends 5-7 context sources and multiple datastores")
    query_rows = [
        [
            "Departure risk and dispatchability (next 90 mins)",
            "AviationWeather METAR/TAF + SIGMET, OpenSky local traffic, OurAirports runway limits, synthetic gate milestones, ASRS crosswind narratives, AD constraints",
            "Microsoft Fabric Eventhouse + Microsoft Fabric Warehouse + Microsoft Fabric Graph (preview) + Azure AI Search",
            "Planner -> parallel retrievers (KQL/SQL/Graph/Vector) -> verifier conflict check -> composed brief",
        ],
        [
            "A321 route change under weather + maintenance pressure",
            "TAF trend deltas, NOTAM impacts, EASA AD text, MEL/tech-log events, alternate airport graph neighbors, historic NTSB analogs",
            "Microsoft Fabric Eventhouse + Microsoft Fabric Warehouse + Microsoft Fabric Graph (preview) + Azure AI Search + Microsoft Fabric Lakehouse",
            "Graph-assisted routing first, then vector evidence, then SQL feasibility scoring",
        ],
        [
            "Disruption chain: crew legality + gate congestion + passenger comms",
            "Crew roster legality tables, turnaround telemetry, baggage disruption patterns, weather hazard nowcast, ASRS station narrative, route network dependencies",
            "Microsoft Fabric Warehouse + Microsoft Fabric Eventhouse + Azure AI Search + Microsoft Fabric Graph (preview)",
            "Multi-agent arbitration: legality agent + ops agent + narrative agent + final verifier",
        ],
    ]
    add_table(
        s13,
        ["Scenario", "Context Sources Used", "Datastores", "Agentic Retrieval Pattern"],
        query_rows,
        [2.2, 4.2, 3.2, 2.98],
        top=1.6,
        height=5.75,
        font_size=8.8,
    )
    add_footer(s13)

    # Slide 14-18
    s14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s14)
    add_store_detail_slide(
        s14,
        "Microsoft Fabric Eventhouse (KQL)",
        [
            "OpenSky states/tracks snapshots.",
            "AviationWeather hazard events.",
            "Recent synthetic ops telemetry.",
            "Windowed anomaly flags for briefing.",
        ],
        [
            "Eventstream or ingestion jobs append events.",
            "KQL retrieves last 30-60 minute context.",
            "Agent router uses time-window predicates.",
            "Verifier checks recency and stale signals.",
        ],
        [
            "Supports disruption-aware brief updates.",
            "Low-latency visibility into rapidly changing ops.",
            "Critical for weather + airport flow decisions.",
            "Improves trust with timestamped evidence.",
        ],
    )

    s15 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s15)
    add_store_detail_slide(
        s15,
        "Microsoft Fabric Warehouse",
        [
            "Airport/runway/navaid dimensions.",
            "Synthetic crew, baggage, gate, MEL facts.",
            "Deterministic schedule and KPI tables.",
            "Curated relational facts for dispatch logic.",
        ],
        [
            "SQL joins integrate ops constraints.",
            "Agent uses rule checks and threshold logic.",
            "Supports explainable KPI outputs.",
            "Feeds scoring for recommended actions.",
        ],
        [
            "Airline decisions need deterministic math.",
            "Best fit for auditable business rules.",
            "Enables repeatable reporting and governance.",
            "Acts as source of truth for structured facts.",
        ],
    )

    s16 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s16)
    add_store_detail_slide(
        s16,
        "Microsoft Fabric Lakehouse",
        [
            "Raw source dumps (JSON/CSV/XML/PDF/ZIP).",
            "Bronze/silver curated Delta layers.",
            "NTSB intermediate transforms from MDB files.",
            "Historical backfill and replay-ready snapshots.",
        ],
        [
            "Ingestion lands raw for lineage.",
            "Spark jobs curate and normalize.",
            "Downstream stores consume curated outputs.",
            "Enables re-indexing and reprocessing.",
        ],
        [
            "Required for evidence traceability.",
            "Prevents lock-in to one retrieval surface.",
            "Supports batch-heavy regulatory pipelines.",
            "Improves resilience for model/data drift fixes.",
        ],
    )

    s17 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s17)
    add_store_detail_slide(
        s17,
        "Azure AI Search (Vector + Hybrid)",
        [
            "Index: idx_ops_narratives (ASRS/NTSB chunks).",
            "Index: idx_regulatory (EASA AD corpus).",
            "Index: idx_airport_ops_docs (ops docs/SOPs).",
            "Chunk metadata: source/date/airport/fleet tags.",
        ],
        [
            "Vector similarity for semantic recall.",
            "Hybrid lexical + semantic for precision.",
            "Metadata filters by fleet, station, timeframe.",
            "Reranking for high-confidence citations.",
        ],
        [
            "Narratives/regulations are unstructured.",
            "Enables case-based reasoning in briefing.",
            "Improves explainability with source citations.",
            "Essential for policy/compliance-aware answers.",
        ],
    )

    s18 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s18)
    add_store_detail_slide(
        s18,
        "Microsoft Fabric Graph (Preview)",
        [
            "Airport-runway-frequency relationships.",
            "Route and dependency edges.",
            "Synthetic ops relationships (flight-crew-gate).",
            "Graph edge list from ops_graph_edges dataset.",
        ],
        [
            "Graph router selects traversal queries.",
            "K-hop neighborhood expansion for alternates.",
            "Combines with SQL constraints for feasibility.",
            "Adds explainable path traces to final answer.",
        ],
        [
            "Disruptions propagate through relationships.",
            "Multi-hop reasoning improves reroute quality.",
            "Supports richer context than flat joins alone.",
            "Key enabler for advanced agentic retrieval.",
        ],
    )

    # Slide 19
    s19 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s19)
    add_header(s19, "Foundry IQ + Agent Framework Retrieval Flow", "How reasoning is layered on top of deterministic retrieval")
    add_bullets(
        s19,
        0.75,
        1.8,
        12.1,
        4.8,
        [
            "Planner Agent (Agent Framework): decomposes query into KQL, SQL, Graph, and vector subtasks.",
            "Retriever Agents: execute in parallel against Microsoft Fabric Eventhouse, Microsoft Fabric Warehouse, Microsoft Fabric Graph, and Azure AI Search.",
            "Verifier Agent: validates freshness, conflict resolution, and regulatory consistency before answer composition.",
            "Foundry IQ orchestration: tracks tool use, context merge strategy, and citation completeness per answer.",
            "Output contract: final brief includes datastore trace, confidence tier, and unresolved-risk checklist.",
        ],
        font_size=16,
    )
    add_footer(s19)

    # Slide 20
    s20 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s20)
    add_header(s20, "Implementation Phases and Success Metrics", "Practical path to production-grade context engineering")
    add_bullets(
        s20,
        0.78,
        1.8,
        12.0,
        4.8,
        [
            "Phase 1 (Weeks 1-2): finalize ingestion contracts and baseline retrieval from each datastore.",
            "Phase 2 (Weeks 3-5): activate multi-agent routing, citation checks, and fallback policies.",
            "Phase 3 (Weeks 6-8): production hardening, RBAC, observability, and cost envelopes.",
            "Primary metrics: citation coverage, retrieval precision@k, median response latency, contradiction rate.",
            "Acceptance threshold: >90% evidence-backed sections and deterministic source trace in every brief.",
        ],
        font_size=16,
    )
    add_footer(s20)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Created: {OUT_PATH}")


if __name__ == "__main__":
    build()
