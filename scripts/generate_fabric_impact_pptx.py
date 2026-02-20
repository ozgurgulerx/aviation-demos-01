#!/usr/bin/env python3
"""
Create a high-impact, non-duplicative deck for the Pilot Brief Fabric architecture.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/THY_Fabric_PilotBrief_Impact_Deck.pptx")

RED = RGBColor(198, 12, 48)
DARK = RGBColor(30, 35, 45)
MID = RGBColor(95, 105, 120)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(0, 145, 90)
AMBER = RGBColor(198, 134, 0)


def bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.24))
    top.fill.solid()
    top.fill.fore_color.rgb = RED
    top.line.fill.background()


def header(slide, title: str, subtitle: str = ""):
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.36), Inches(12.1), Inches(0.82)).text_frame
    tb.clear()
    p = tb.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Segoe UI"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = DARK
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.68), Inches(1.02), Inches(12.0), Inches(0.45)).text_frame
        sb.clear()
        sp = sb.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Segoe UI"
        sr.font.size = Pt(13)
        sr.font.color.rgb = MID


def footer(slide):
    fb = slide.shapes.add_textbox(Inches(0.65), Inches(7.0), Inches(12.1), Inches(0.28)).text_frame
    fb.clear()
    p = fb.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "Turkish Airlines Pilot Brief Demo - Fabric-first architecture"
    r.font.name = "Segoe UI"
    r.font.size = Pt(10)
    r.font.color.rgb = MID


def bullets(slide, items, x=0.8, y=1.7, w=12.0, h=5.2, size=18):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK


def table(slide, headers, rows, col_widths, top=1.62, height=5.5, font_size=9.5):
    t = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.38), Inches(top), Inches(12.58), Inches(height)).table
    for i, w in enumerate(col_widths):
        t.columns[i].width = Inches(w)
    for c, h in enumerate(headers):
        cell = t.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
    for r_i, row in enumerate(rows, start=1):
        for c_i, val in enumerate(row):
            cell = t.cell(r_i, c_i)
            cell.text = val
            if r_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(234, 239, 245)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK
    return t


def metric_card(slide, x, y, w, h, title, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(1.7)
    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = MID
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = DARK


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1) Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Pilot Brief Generator - Architecture and Data Story", "High-value view: what is deployed, why this structure, and where it goes next")
    bullets(
        s,
        [
            "Objective: maximize retrieval quality, traceability, and operational credibility.",
            "Principle: one primary datastore per domain in target architecture.",
            "Scope: Fabric Lakehouse + KQL + Warehouse + Graph + Azure AI Search.",
        ],
        size=19,
    )
    footer(s)

    # 2) Executive snapshot
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Executive Snapshot", "Current deployed footprint and readiness")
    metric_card(s, 0.8, 1.8, 2.9, 1.35, "Lakehouse files", "132", GREEN)
    metric_card(s, 4.05, 1.8, 2.9, 1.35, "Lakehouse size", "871.7 MB", GREEN)
    metric_card(s, 7.3, 1.8, 2.9, 1.35, "KQL tables", "6", GREEN)
    metric_card(s, 10.55, 1.8, 2.0, 1.35, "KQL rows", "13.6K+", GREEN)
    bullets(
        s,
        [
            "Fabric landing and event analytics are live now (Lakehouse + Eventhouse/KQL).",
            "Vector/hybrid retrieval is active in Azure AI Search for narrative and regulatory corpora.",
            "Target architecture avoids permanent secondary stores; current duplicates are migration artifacts only.",
        ],
        y=3.55,
        h=2.7,
        size=17,
    )
    footer(s)

    # 3) Data and purpose
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Data and Purpose (Detailed)", "Aviation relevance of each source family")
    headers = ["Dataset / Source", "Format", "Why Important for Aviation", "Pilot Brief Value"]
    rows = [
        ["OpenSky state vectors", "JSON snapshots", "Live aircraft context and airspace activity awareness", "Real-time grounding for disruptions and traffic context"],
        ["AviationWeather hazards", "CSV/XML/JSON.gz", "Operational weather risk and route hazard visibility", "Safety-centric risk flags in pre-flight brief"],
        ["NOTAM exports", "JSON/JSONL", "Mandatory restrictions and procedural constraints", "Compliance-aware briefing content"],
        ["OurAirports + OpenFlights", "CSV + DAT", "Airport/runway/nav reference backbone", "Deterministic joins for airport feasibility context"],
        ["ASRS narratives", "JSONL", "Crew-reported hazards and recurring human-factor patterns", "Narrative analog retrieval for risk explanation"],
        ["NTSB archives", "ZIP/MDB", "Historical incidents and structured causal context", "Long-horizon enrichment of threat patterns"],
        ["EASA AD corpus", "PDF + TSV", "Airworthiness and regulatory directive retrieval", "Deterministic + semantic regulatory evidence"],
        ["Synthetic ops overlay", "CSV", "Crew/gate/baggage/MEL realism not openly available", "End-to-end operational scenario simulation"],
    ]
    table(s, headers, rows, [3.0, 1.7, 4.05, 3.83], font_size=8.8)
    footer(s)

    # 4) Datastore role model
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Workload-by-Datastore Model", "The deliberate split")
    bullets(
        s,
        [
            "KQL/Eventhouse: live telemetry + recent event context.",
            "Warehouse/SQL: stable relational joins, KPIs, deterministic reporting.",
            "Lakehouse: raw landing + batch curation.",
            "AI Search: vector/hybrid retrieval for narratives/docs.",
            "Graph: relationship routing and multi-hop context.",
        ],
        size=20,
    )
    footer(s)

    # 5) Single-store target mapping with reasoning
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Single-Store Mapping (Target)", "Primary datastore selection and reasoning")
    headers = ["Domain", "Primary Datastore", "Retrieval Mode", "Reasoning"]
    rows = [
        ["Real-time ops/weather", "Fabric KQL/Eventhouse", "Time-window and anomaly lookup", "Best latency/throughput fit for append-heavy event streams"],
        ["Airport/reference dimensions", "Fabric Warehouse/SQL", "Deterministic joins and KPI queries", "Relational integrity and predictable query plans"],
        ["Narratives and regulatory corpus", "Azure AI Search", "Hybrid vector + keyword", "Highest retrieval quality on long unstructured text"],
        ["Graph context", "Fabric Graph model", "Path and neighborhood traversal", "Explicit multi-hop relationships for routing and explainability"],
        ["Raw archives and replay", "Fabric Lakehouse", "Batch transforms and reprocessing", "Low-cost durable store with strong ingestion flexibility"],
        ["Application state", "PostgreSQL", "Transactional API workloads", "Operationally simple runtime persistence for app state"],
    ]
    table(s, headers, rows, [2.25, 2.95, 2.75, 4.55], font_size=9.2)
    footer(s)

    # 6) Query routing examples
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Agentic Retrieval Routing", "How different queries trigger different stores")
    headers = ["Example Query", "Primary Route", "Secondary Step", "Why This Path Wins"]
    rows = [
        ["What changed in last 45 minutes around LTFM?", "KQL", "Warehouse for static airport context", "Freshness first, then deterministic enrichment"],
        ["Any active restrictions affecting alternate airports?", "AI Search (NOTAM/AD)", "Warehouse join to airports/runways", "Regulatory text recall plus precise filter"],
        ["Why is this scenario riskier than normal?", "AI Search (ASRS/NTSB)", "KQL for current hazards", "Narrative analogs validated against live signals"],
        ["Which entities are linked to this disruption chain?", "Graph", "KQL timeline overlay", "Multi-hop explainability plus event chronology"],
    ]
    table(s, headers, rows, [3.45, 2.05, 2.35, 4.73], font_size=9.0)
    footer(s)

    # 7) Why this works
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Why This Structure Delivers Impact", "Architecture rationale")
    bullets(
        s,
        [
            "Latency fit: live questions hit KQL, not batch stores.",
            "Correctness fit: deterministic metrics stay in SQL/Warehouse.",
            "Recall fit: long-form documents use vector + hybrid retrieval.",
            "Traceability fit: each answer carries source-level evidence.",
            "Scalability fit: ingestion, analytics, and retrieval scale independently.",
            "Resilience fit: subsystem failure degrades gracefully instead of failing hard.",
        ],
        size=18,
    )
    footer(s)

    # 8) Gaps and plan
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Gaps and 30-60-90 Plan", "Execution-focused next steps")
    bullets(
        s,
        [
            "30 days: materialize Lakehouse bronze/silver tables and publish Warehouse core facts/dims.",
            "60 days: promote Fabric graph model and tighten query router policies with benchmark suite.",
            "90 days: production hardening (SLAs, lineage, observability, governance controls).",
            "Success criteria: lower hallucination risk, faster response, higher source traceability coverage.",
        ],
        size=18,
    )
    footer(s)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
