#!/usr/bin/env python3
"""Generate a PPTX explaining the Fabric Intent Graph and its role in query routing."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/Intent_Graph_Query_Routing.pptx")

RED = RGBColor(198, 12, 48)
DARK = RGBColor(30, 35, 45)
MID = RGBColor(95, 105, 120)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(0, 145, 90)
AMBER = RGBColor(198, 134, 0)
BLUE = RGBColor(0, 100, 180)
PANEL = RGBColor(233, 238, 247)


def bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.24))
    top.fill.solid()
    top.fill.fore_color.rgb = RED
    top.line.fill.background()


def header(slide, title: str, subtitle: str = ""):
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.36), Inches(12.1), Inches(0.82)).text_frame
    tb.clear()
    p = tb.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Segoe UI"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = DARK
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.68), Inches(1.02), Inches(12.0), Inches(0.45)).text_frame
        sb.clear()
        sp = sb.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Segoe UI"
        sr.font.size = Pt(13)
        sr.font.color.rgb = MID


def footer(slide):
    fb = slide.shapes.add_textbox(Inches(0.65), Inches(7.0), Inches(12.1), Inches(0.28)).text_frame
    fb.clear()
    p = fb.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "Aviation RAG - Intent Graph & Query Routing"
    r.font.name = "Segoe UI"
    r.font.size = Pt(10)
    r.font.color.rgb = MID


def bullets(slide, items, x=0.8, y=1.7, w=12.0, h=5.2, size=18):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK


def add_table(slide, headers, rows, col_widths, top=1.62, height=5.5, font_size=9.5):
    t = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.38), Inches(top), Inches(12.58), Inches(height)).table
    for i, w in enumerate(col_widths):
        t.columns[i].width = Inches(w)
    for c, h in enumerate(headers):
        cell = t.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
    for r_i, row in enumerate(rows, start=1):
        for c_i, val in enumerate(row):
            cell = t.cell(r_i, c_i)
            cell.text = val
            if r_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(234, 239, 245)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK
    return t


def metric_card(slide, x, y, w, h, title, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(1.7)
    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = MID
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = DARK


def labeled_box(slide, x, y, w, h, label, body, border_color, body_size=11):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = label
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = border_color
    p1.space_after = Pt(4)
    for line in body:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(body_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(2)


def arrow_label(slide, x, y, text, color=MID):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.0), Inches(0.3)).text_frame
    tb.clear()
    p = tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = color


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Intent Graph & Query Routing", "How the Fabric-hosted intent graph drives evidence-aware retrieval planning")
    bullets(s, [
        "The intent graph is the declarative contract between user intents, evidence types, and data tools.",
        "It replaces hardcoded routing rules with a structured, externally manageable knowledge model.",
        "The graph can be hosted in Fabric (primary), loaded from JSON (secondary), or use a built-in default.",
        "This deck explains the graph structure, how it drives query routing, and end-to-end execution flow.",
    ], size=17)
    footer(s)

    # ── Slide 2: Graph Structure Overview ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Intent Graph Structure", "Three entity types connected by three edge types")

    metric_card(s, 0.8, 1.6, 2.6, 1.2, "Intents", "5", RED)
    metric_card(s, 3.7, 1.6, 2.6, 1.2, "Evidence Types", "6", BLUE)
    metric_card(s, 6.6, 1.6, 2.6, 1.2, "Tools", "7", GREEN)
    metric_card(s, 9.5, 1.6, 3.2, 1.2, "Edge Types", "3", AMBER)

    # Entity listings as labeled boxes
    labeled_box(s, 0.8, 3.2, 3.5, 3.2, "Intents (what user wants)", [
        "PilotBrief.Departure",
        "PilotBrief.Arrival",
        "Disruption.Explain",
        "Policy.Check",
        "Replay.History",
    ], RED)

    labeled_box(s, 4.65, 3.2, 3.5, 3.2, "Evidence (what is needed)", [
        "METAR - weather obs",
        "TAF - weather forecast",
        "NOTAM - restrictions",
        "RunwayConstraints",
        "Hazards - risk signals",
        "SOPClause - procedures",
    ], BLUE)

    labeled_box(s, 8.5, 3.2, 4.2, 3.2, "Tools (where to get it)", [
        "GRAPH - relationship traversal",
        "KQL - live event windows",
        "SQL - relational joins/KPIs",
        "NOSQL - document lookups",
        "VECTOR_REG - regulatory corpus",
        "VECTOR_OPS - ops narratives",
        "VECTOR_AIRPORT - airport docs",
    ], GREEN)

    footer(s)

    # ── Slide 3: Edge Types ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Graph Edges: The Routing Contract", "Three edge types define intent-to-tool mapping")

    add_table(s, ["Edge Type", "From", "To", "Key Attribute", "Purpose"], [
        ["requires", "Intent", "Evidence", "optional: true/false", "Defines what evidence each intent needs. Mandatory evidence must be fulfilled for verified answers."],
        ["authoritative_in", "Evidence", "Tool", "priority: 1, 2, ...", "Maps each evidence type to the best tool to retrieve it, ranked by priority."],
        ["expansion_rules", "Intent", "Tool", "reason: text", "Triggers entity expansion (e.g. airport -> stations/alternates) before evidence collection."],
    ], [1.6, 1.1, 1.1, 1.9, 6.88], font_size=10.0, top=1.55, height=2.8)

    bullets(s, [
        "Example: PilotBrief.Departure requires METAR (mandatory). METAR is authoritative_in KQL (priority 1).",
        "Result: the planner schedules a KQL lookup for METAR evidence when handling departure briefing queries.",
        "If KQL is unavailable, the system falls back to the next priority tool for that evidence type.",
        "expansion_rules trigger a GRAPH call first for PilotBrief intents to resolve airport -> stations -> alternates.",
    ], y=4.6, h=2.5, size=15)
    footer(s)

    # ── Slide 4: Requires Edges Detail ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "'requires' Edges: Intent-Evidence Matrix", "Which evidence each intent needs (mandatory vs optional)")

    add_table(s, ["Intent", "METAR", "TAF", "NOTAM", "Runway\nConstraints", "Hazards", "SOPClause"], [
        ["PilotBrief.Departure", "Required", "Required", "Required", "Required", "Optional", "-"],
        ["PilotBrief.Arrival", "Required", "Required", "Required", "Required", "-", "-"],
        ["Disruption.Explain", "Optional", "-", "Required", "-", "Required", "Optional"],
        ["Policy.Check", "-", "-", "-", "-", "-", "Required"],
        ["Replay.History", "Required", "-", "Optional", "-", "Required", "-"],
    ], [2.4, 1.65, 1.65, 1.65, 1.65, 1.65, 1.65], font_size=10.0, top=1.55, height=3.5)

    bullets(s, [
        "Mandatory evidence must be retrieved for the answer to be marked 'verified'.",
        "Optional evidence enriches the response but its absence does not block verification.",
        "The coverage checklist in the retrieval plan tracks planned/missing status per evidence slot.",
    ], y=5.3, h=1.6, size=14)
    footer(s)

    # ── Slide 5: Authoritative-In Edges ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "'authoritative_in' Edges: Evidence-Tool Mapping", "Which tool is the best source for each evidence type")

    add_table(s, ["Evidence Type", "Priority 1 (Primary)", "Priority 2 (Fallback)", "Rationale"], [
        ["METAR", "KQL (Eventhouse)", "-", "Live weather observations are time-series data, best served by KQL windows"],
        ["TAF", "KQL (Eventhouse)", "-", "Forecast data follows same event-stream pattern as METAR"],
        ["NOTAM", "NOSQL", "VECTOR_REG", "Primary: structured document store. Fallback: semantic search over regulatory corpus"],
        ["RunwayConstraints", "SQL (Warehouse)", "-", "Deterministic relational data: runway lengths, closures, configurations"],
        ["Hazards", "KQL (Eventhouse)", "-", "Real-time hazard signals from telemetry streams"],
        ["SOPClause", "VECTOR_REG", "-", "Long-form regulatory/procedure text needs vector/hybrid retrieval"],
    ], [2.0, 2.5, 2.5, 5.58], font_size=9.5, top=1.55, height=4.3)

    bullets(s, [
        "Priority ordering enables graceful degradation: if the primary tool is down, the next-priority tool is used.",
    ], y=6.1, h=0.8, size=14)
    footer(s)

    # ── Slide 6: Agentic Routing Flow ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Agentic Routing: End-to-End Flow", "How a user query becomes a multi-source retrieval plan")

    # Flow boxes
    labeled_box(s, 0.5, 1.6, 2.5, 1.4, "1. User Query", [
        "Natural language",
        "question from pilot",
        "or ops analyst",
    ], RED, body_size=10)

    arrow_label(s, 2.7, 2.1, "-->")

    labeled_box(s, 3.2, 1.6, 2.8, 1.4, "2. Load Intent Graph", [
        "Fabric endpoint (primary)",
        "JSON file (secondary)",
        "Builtin default (fallback)",
    ], BLUE, body_size=10)

    arrow_label(s, 5.7, 2.1, "-->")

    labeled_box(s, 6.3, 1.6, 3.2, 1.4, "3. LLM Router (gpt-5-nano)", [
        "Receives: query + graph +",
        "schemas + tool catalog",
        "Outputs: AgenticPlan JSON",
    ], AMBER, body_size=10)

    arrow_label(s, 9.2, 2.1, "-->")

    labeled_box(s, 9.8, 1.6, 3.0, 1.4, "4. Post-Validation", [
        "Graph walk verifies all",
        "required evidence has",
        "a scheduled tool call",
    ], GREEN, body_size=10)

    # Second row
    labeled_box(s, 0.5, 3.5, 3.5, 1.6, "5. Plan Execution", [
        "ThreadPoolExecutor runs tool",
        "calls in parallel (respecting",
        "depends_on ordering)",
        "GRAPH expansion runs first",
    ], RED, body_size=10)

    arrow_label(s, 3.7, 4.1, "-->")

    labeled_box(s, 4.3, 3.5, 3.5, 1.6, "6. Evidence Verification", [
        "Coverage checklist: each",
        "required evidence slot is",
        "checked against actual results",
        "Status: planned | filled | missing",
    ], BLUE, body_size=10)

    arrow_label(s, 7.5, 4.1, "-->")

    labeled_box(s, 8.1, 3.5, 4.7, 1.6, "7. Context Reconciliation & Synthesis", [
        "Source results reconciled, conflicts detected",
        "Context text composed and sent to LLM",
        "Answer synthesized with source-level citations",
        "is_verified = all required evidence filled",
    ], GREEN, body_size=10)

    bullets(s, [
        "The AgenticPlan JSON includes: intent classification, entity extraction, required_evidence list, ordered tool_calls with dependencies, and a coverage checklist.",
        "If the LLM router fails, a deterministic fallback walks the same graph edges without an LLM call.",
        "The legacy retrieval planner (keyword heuristics) is a second fallback that does not use the intent graph.",
    ], y=5.5, h=1.6, size=13)
    footer(s)

    # ── Slide 7: Concrete Routing Example ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Worked Example: Departure Briefing Query", "\"What are the current conditions and restrictions for departure from LTFM?\"")

    add_table(s, ["Step", "Action", "Graph Edge Used", "Result"], [
        ["1. Intent Classification", "LLM classifies as PilotBrief.Departure", "-", "intent = PilotBrief.Departure, confidence = 0.92"],
        ["2. Evidence Requirements", "Graph walk: requires edges for PilotBrief.Departure", "requires", "METAR(req), TAF(req), NOTAM(req), RunwayConstraints(req), Hazards(opt)"],
        ["3. Entity Expansion", "Graph traversal: LTFM -> stations, alternates", "expansion_rules", "GRAPH call scheduled first (call_1), depends_on = []"],
        ["4. METAR Retrieval", "authoritative_in: METAR -> KQL (priority 1)", "authoritative_in", "KQL call for METAR, depends_on = [call_1]"],
        ["5. TAF Retrieval", "authoritative_in: TAF -> KQL (priority 1)", "authoritative_in", "KQL call for TAF, depends_on = [call_1]"],
        ["6. NOTAM Retrieval", "authoritative_in: NOTAM -> NOSQL (p1), VECTOR_REG (p2)", "authoritative_in", "NOSQL call for NOTAM, depends_on = [call_1]"],
        ["7. RunwayConstraints", "authoritative_in: RunwayConstraints -> SQL (priority 1)", "authoritative_in", "SQL call for runway data, depends_on = [call_1]"],
        ["8. Hazards (optional)", "authoritative_in: Hazards -> KQL (priority 1)", "authoritative_in", "KQL call for hazards, depends_on = [call_1]"],
        ["9. Coverage Check", "Verify all 4 required + 1 optional evidence slots", "-", "is_verified = true if all required slots filled"],
    ], [1.5, 3.6, 1.8, 5.68], font_size=8.5, top=1.55, height=5.2)

    footer(s)

    # ── Slide 8: Three-Tier Fallback ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Three-Tier Graph Loading", "Resilient loading strategy ensures routing always works")

    labeled_box(s, 0.8, 1.7, 3.6, 2.8, "Tier 1: Fabric Graph Endpoint", [
        "POST to FABRIC_GRAPH_ENDPOINT",
        "Bearer token authentication",
        "Returns live intent graph snapshot",
        "Allows external graph management",
        "without code deployments",
        "",
        "Status: blocked (FDPO capacity",
        "assignment pending)",
    ], RED, body_size=10)

    labeled_box(s, 4.8, 1.7, 3.6, 2.8, "Tier 2: JSON File", [
        "Reads from INTENT_GRAPH_JSON_PATH",
        "or data/intent_graph.json",
        "Useful for local dev and testing",
        "Can be updated by CI/CD pipeline",
        "",
        "",
        "Status: available as override",
        "",
    ], AMBER, body_size=10)

    labeled_box(s, 8.8, 1.7, 3.8, 2.8, "Tier 3: Built-in Default", [
        "Hardcoded DEFAULT_INTENT_GRAPH",
        "in intent_graph_provider.py",
        "Always available, zero dependencies",
        "5 intents, 6 evidence types, 7 tools",
        "17 requires edges, 7 authoritative_in",
        "2 expansion_rules",
        "",
        "Status: currently active",
    ], GREEN, body_size=10)

    bullets(s, [
        "Loading priority: Fabric endpoint > JSON file > built-in default. First successful load wins.",
        "The IntentGraphSnapshot.source field tracks which tier was used (fabric-graph | json-file | builtin-default).",
        "This source label is included in retrieval plan reasoning for full observability.",
        "When Fabric capacity assignment completes, Tier 1 will enable graph updates without redeployment.",
    ], y=4.8, h=2.1, size=14)
    footer(s)

    # ── Slide 9: Agentic vs Legacy Routing ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Agentic vs Legacy Routing", "Two retrieval planning paths with different graph dependency")

    add_table(s, ["Dimension", "Agentic Path (primary)", "Legacy Path (fallback)"], [
        ["Graph Dependency", "Full intent graph sent to LLM + post-validation walk", "No graph dependency; keyword heuristics only"],
        ["Intent Classification", "LLM-based (gpt-5-nano), high confidence", "Keyword matching (_infer_intent)"],
        ["Evidence Planning", "Graph-driven: requires + authoritative_in edges", "Profile-driven: pilot-brief, compliance, etc."],
        ["Tool Selection", "Per-evidence priority from authoritative_in", "Source activation by keyword signals"],
        ["Coverage Verification", "Full evidence slot tracking with is_verified", "Citation count only"],
        ["Entity Expansion", "GRAPH tool call via expansion_rules", "Not supported"],
        ["Fallback Trigger", "Default path when agentic_enabled=true", "Used when orchestrator init fails or LLM unavailable"],
        ["Route Label", "AGENTIC", "SQL / SEMANTIC / HYBRID"],
    ], [2.0, 5.29, 5.29], font_size=9.2, top=1.55, height=5.0)

    footer(s)

    # ── Slide 10: Key Takeaways ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Key Takeaways", "Why intent graph routing matters for aviation RAG")

    bullets(s, [
        "Declarative contract: routing rules live in data (the graph), not in code. Changes don't require redeployment.",
        "Evidence completeness: the graph defines what evidence is required per intent, enabling verified answers.",
        "Graceful degradation: priority-ranked tool mappings + three-tier loading means routing always works.",
        "Observability: every retrieval plan includes graph_source, coverage checklist, and tool execution traces.",
        "Separation of concerns: domain experts manage the graph; the planner and executor are generic.",
        "Future-ready: when Fabric hosting is live, graph updates flow to production without code changes.",
    ], size=17)
    footer(s)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
