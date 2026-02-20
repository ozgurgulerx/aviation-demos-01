#!/usr/bin/env python3
"""
Generate a concise service-mapping deck for the aviation demo.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/THY_Fabric_Data_Service_Mapping_Readable.pptx")

COLOR_RED = RGBColor(198, 12, 48)
COLOR_DARK = RGBColor(30, 35, 45)
COLOR_LIGHT = RGBColor(245, 247, 250)
COLOR_MID = RGBColor(99, 110, 126)
COLOR_WHITE = RGBColor(255, 255, 255)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT
    bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
    top.fill.solid()
    top.fill.fore_color.rgb = COLOR_RED
    top.line.fill.background()


def add_header(slide, title: str, subtitle: str = ""):
    tf = slide.shapes.add_textbox(Inches(0.7), Inches(0.38), Inches(12.0), Inches(0.8)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Segoe UI"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = COLOR_DARK

    if subtitle:
        sf = slide.shapes.add_textbox(Inches(0.72), Inches(1.02), Inches(12.0), Inches(0.45)).text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Segoe UI"
        sr.font.size = Pt(13)
        sr.font.color.rgb = COLOR_MID


def add_footer(slide, text="Turkish Airlines Demo - Service Mapping"):
    tf = slide.shapes.add_textbox(Inches(0.7), Inches(7.02), Inches(12.0), Inches(0.3)).text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(10)
    r.font.color.rgb = COLOR_MID


def add_bullets(slide, bullets, x=0.8, y=1.7, w=12.0, h=5.2, size=18):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = COLOR_DARK


def add_table(slide, headers, rows, top=1.65, height=5.35, font_size=10):
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(0.4), Inches(top), Inches(12.55), Inches(height)
    ).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_RED
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE

    for r_i, row in enumerate(rows, start=1):
        for c_i, value in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = value
            if r_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(235, 239, 245)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLOR_DARK

    return table


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s1)
    add_header(s1, "Pilot Brief Demo - Data by Service", "Readable deployment map (current state)")
    add_bullets(
        s1,
        [
            "This version shows exactly where each dataset lives today.",
            "Fabric is now primary for landing + event analytics.",
            "Azure AI Search remains primary for vector/hybrid retrieval.",
            "PostgreSQL remains app/relational fallback until Fabric Warehouse curation is completed.",
            "Target architecture is single authoritative datastore per domain (no permanent secondary store).",
        ],
        size=18,
    )
    add_footer(s1)

    # Slide 2: Fabric now
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s2)
    add_header(s2, "Fabric-Resident Data (Now)", "Lakehouse + Eventhouse/KQL")
    headers = ["Fabric Service", "Dataset Group", "Current Volume", "Purpose in Demo"]
    rows = [
        ["Lakehouse Files", "Full raw landing zone (10 groups)", "132 files, 871,733,208 bytes", "Canonical ingestion source"],
        ["Lakehouse Files", "ASRS processed", "3 files, 630.6 MB", "Large narrative corpus for RAG prep"],
        ["Lakehouse Files", "NTSB archives", "2 zip files, 132.8 MB", "Batch incident history"],
        ["KQL (Eventhouse)", "opensky_states", "10,867 rows", "Live aircraft-state context"],
        ["KQL (Eventhouse)", "hazards_airsigmets", "4 rows", "SIGMET/AIRMET signal"],
        ["KQL (Eventhouse)", "hazards_gairmets", "209 rows", "G-AIRMET hazard signal"],
        ["KQL (Eventhouse)", "hazards_aireps_raw", "2,039 rows", "Raw PIREP/AIREP retrieval"],
        ["KQL (Eventhouse)", "ops_graph_edges", "520 rows", "Operational relationship edges"],
    ]
    table = add_table(s2, headers, rows, font_size=10)
    table.columns[0].width = Inches(2.1)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(2.7)
    table.columns[3].width = Inches(4.75)
    add_footer(s2)

    # Slide 3: Detailed data catalog
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s3)
    add_header(s3, "Data and Purpose (Detailed)", "What we use, why it matters, and aviation value")
    headers = ["Dataset / Source", "Format", "Aviation Purpose", "Demo Value"]
    rows = [
        ["OpenSky state vectors (Istanbul-focused)", "JSON (API snapshots)", "Live aircraft movement context and disruption awareness", "Real-time situational grounding in pilot brief"],
        ["AviationWeather METAR/TAF/SIGMET/G-AIRMET/PIREP", "CSV/XML/JSON.gz", "Weather and hazard risk interpretation", "Operational risk flags and route caution"],
        ["NOTAM (FAA/NASA query exports)", "JSON/JSONL", "Airspace and airport restriction awareness", "Regulatory compliance and briefing constraints"],
        ["OurAirports + OpenFlights", "CSV + DAT", "Airport/runway/nav/reference backbone", "Join keys and route feasibility context"],
        ["ASRS incident narratives", "JSONL (processed)", "Crew-reported safety patterns and lessons learned", "Narrative analog retrieval and risk explanation"],
        ["NTSB accident archives", "ZIP/MDB", "Historical incident structure and causal context", "Long-horizon risk enrichment"],
        ["EASA AD corpus (PDF + metadata)", "PDF + TSV", "Airworthiness and maintenance/regulatory directives", "Deterministic + semantic regulatory retrieval"],
        ["Synthetic ops overlay (crew, gate, baggage, MEL)", "CSV", "Operational realism not publicly available at event level", "End-to-end agentic scenario simulation"],
    ]
    table = add_table(s3, headers, rows, font_size=8.6)
    table.columns[0].width = Inches(3.3)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(3.75)
    table.columns[3].width = Inches(3.7)
    add_footer(s3)

    # Slide 4: Data importance + datastore fitness
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s4)
    add_header(s4, "Data Importance and Datastore Fit", "Why each dataset matters in aviation")
    headers = ["Data Domain", "Why Important for Aviation", "Primary Datastore", "What This Store Is Good For"]
    rows = [
        ["OpenSky + weather hazards", "Pre-flight and in-flight situational awareness", "Fabric KQL/Eventhouse", "Low-latency event/time-window queries"],
        ["NOTAM + AD + safety notices", "Regulatory and operational restrictions", "Azure AI Search", "Hybrid semantic retrieval over long text"],
        ["Airport/runway/nav references", "Route feasibility and airport suitability", "Warehouse/SQL", "Reliable joins and deterministic filters"],
        ["ASRS/NTSB narratives", "Risk patterns and incident analogs", "Azure AI Search", "Vector recall for narrative similarity"],
        ["Synthetic ops overlays", "Demo realism for crew/gate/baggage/MEL use-cases", "Warehouse/SQL + KQL", "Blend of analytics + recent event context"],
        ["Raw archives and source evidence", "Auditability and reprocessing", "Fabric Lakehouse", "Cheap durable landing + batch transformation"],
        ["Entity relationships (airport->flight->hazard)", "Explainable multi-hop retrieval", "Graph model", "Path-based reasoning and context routing"],
    ]
    table = add_table(s4, headers, rows, font_size=9)
    table.columns[0].width = Inches(2.55)
    table.columns[1].width = Inches(4.05)
    table.columns[2].width = Inches(2.55)
    table.columns[3].width = Inches(3.40)
    add_footer(s4)

    # Slide 5: Cross-service map with reasoning (single primary store target)
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s5)
    add_header(s5, "Cross-Service Mapping (Target)", "Single primary store per domain")
    headers = ["Domain", "Primary Store", "Retrieval Mode", "Reasoning"]
    rows = [
        ["Real-time ops/weather", "Fabric KQL/Eventhouse", "Time-window + anomaly lookup", "Streaming/event-native, very fast recent-window filters"],
        ["Airport/reference dimensions", "Fabric Warehouse/SQL", "Deterministic SQL joins", "Relational integrity and predictable KPI logic"],
        ["Unstructured safety/regulatory", "Azure AI Search", "Hybrid vector + keyword", "Best recall/precision tradeoff for long document retrieval"],
        ["Graph context", "Fabric Graph model", "Graph-assisted routing", "Multi-hop relationship traversal and explainable paths"],
        ["Synthetic operations overlay", "Fabric Warehouse + KQL split", "Scenario joins + event context", "Static facts in SQL, recent signals in KQL"],
        ["Application transactional state", "PostgreSQL", "App runtime queries", "Operational simplicity and existing backend compatibility"],
    ]
    table = add_table(s5, headers, rows, font_size=9)
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.55)
    table.columns[3].width = Inches(5.0)
    add_footer(s5)

    # Slide 6: Workload pattern
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s6)
    add_header(s6, "Workload-by-Service Pattern", "Deliberate platform split")
    add_bullets(
        s6,
        [
            "KQL/Eventhouse: live telemetry + recent event context.",
            "Warehouse/SQL: stable relational joins, KPIs, deterministic reporting.",
            "Lakehouse: raw landing + batch curation.",
            "AI Search: vector/hybrid retrieval for narratives/docs.",
            "Graph: relationship routing/multi-hop context.",
        ],
        size=18,
    )
    add_footer(s6)

    # Slide 7: Why this structure
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s7)
    add_header(s7, "Why We Use This Structure", "Reasoning behind the architecture")
    add_bullets(
        s7,
        [
            "Latency fit: real-time questions hit KQL, not batch stores.",
            "Correctness fit: deterministic calculations stay in SQL/Warehouse.",
            "Recall fit: unstructured docs use vector + hybrid retrieval in AI Search.",
            "Traceability fit: each answer is backed by explicit source evidence.",
            "Scalability fit: ingestion, analytics, and semantic retrieval scale independently.",
            "Resilience fit: one subsystem outage does not collapse the entire retrieval path.",
        ],
        size=18,
    )
    add_footer(s7)

    # Slide 8: Immediate execution
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s8)
    add_header(s8, "Next 5 Actions", "To make Fabric even more representative")
    add_bullets(
        s8,
        [
            "1) Materialize Lakehouse bronze/silver tables from ingest_full.",
            "2) Publish curated Fabric Warehouse facts/dimensions for SQL serving.",
            "3) Promote graph_nodes + graph_edges model in Fabric (preview path).",
            "4) Re-index AI Search from curated Fabric outputs with consistent metadata.",
            "5) Record benchmark prompts with source-attribution screenshots for leadership demo.",
        ],
        size=18,
    )
    add_footer(s8)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
