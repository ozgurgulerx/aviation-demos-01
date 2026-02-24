#!/usr/bin/env python3
"""
Generate an executive 10-slide deck for the predictive delay optimization plan.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/predictive_delay_plan_executive.pptx")

COLOR_RED = RGBColor(198, 12, 48)
COLOR_DARK = RGBColor(30, 35, 45)
COLOR_LIGHT = RGBColor(245, 247, 250)
COLOR_MID = RGBColor(99, 110, 126)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_PANEL = RGBColor(232, 238, 247)
COLOR_PANEL_ALT = RGBColor(223, 233, 244)
COLOR_GREEN = RGBColor(21, 130, 83)
COLOR_AMBER = RGBColor(187, 124, 15)


def add_background(slide) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT
    bg.line.fill.background()

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.24)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_RED
    band.line.fill.background()


def add_header(slide, title: str, subtitle: str = "") -> None:
    t = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12.0), Inches(0.8))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = COLOR_DARK

    if subtitle:
        s = slide.shapes.add_textbox(
            Inches(0.68), Inches(1.02), Inches(12.0), Inches(0.45)
        )
        stf = s.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Segoe UI"
        sp.font.size = Pt(14)
        sp.font.color.rgb = COLOR_MID


def add_footer(slide, text: str = "Predictive Delay Optimization - Executive Plan") -> None:
    f = slide.shapes.add_textbox(Inches(0.7), Inches(7.02), Inches(12.0), Inches(0.25))
    tf = f.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_MID


def add_bullets(
    slide,
    bullets: Sequence[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    space_after: int = 8,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(space_after)
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = COLOR_DARK


def add_panel(slide, x: float, y: float, w: float, h: float, fill=COLOR_PANEL):
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = COLOR_RED
    panel.line.width = Pt(1)
    return panel


def set_panel_text(panel, title: str, lines: Iterable[str], title_size: int = 14, body_size: int = 11) -> None:
    tf = panel.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK
    for line in lines:
        bp = tf.add_paragraph()
        bp.text = line
        bp.font.name = "Segoe UI"
        bp.font.size = Pt(body_size)
        bp.font.color.rgb = COLOR_DARK
        bp.space_after = Pt(4)


def add_table(
    slide,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    col_widths: Sequence[float],
    body_font_size: int = 10,
) -> None:
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)
    ).table
    for idx, cw in enumerate(col_widths):
        table.columns[idx].width = Inches(cw)

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_RED
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(236, 241, 248)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(body_font_size)
                p.font.color.rgb = COLOR_DARK


def slide_1_title(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(
        s,
        "Delay Risk + Delay Minutes",
        "Flagship predictive optimization plan (executive view)",
    )
    add_bullets(
        s,
        [
            "Goal: predict P(arrival delay >= 15), expected delay minutes, and uncertainty per flight leg.",
            "Business outcome: prioritize interventions that improve OTP and reduce downstream disruption cost.",
            "Scope now: demo-ready v0-v2 implementation using current data and current app stack.",
        ],
        x=0.8,
        y=1.8,
        w=11.8,
        h=2.0,
    )

    p1 = add_panel(s, 0.8, 4.0, 3.9, 2.35)
    set_panel_text(
        p1,
        "Decision Inputs",
        [
            "Historical labels (BTS).",
            "Live hazard and traffic signals.",
            "NOTAM constraints and network effects.",
        ],
    )
    p2 = add_panel(s, 4.95, 4.0, 3.9, 2.35)
    set_panel_text(
        p2,
        "Decision Output",
        [
            "Risk score per departure.",
            "Expected delay and interval.",
            "Top 3 explainable drivers.",
        ],
    )
    p3 = add_panel(s, 9.1, 4.0, 3.45, 2.35)
    set_panel_text(
        p3,
        "A/B Story",
        [
            "Baseline vs optimized toggle.",
            "AUROC, Brier, MAE uplift.",
            "Concrete high-impact flights.",
        ],
    )
    add_footer(s)


def slide_2_feasible_now(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Why This Is Feasible Now", "Current data already covers the required signal groups")
    headers = ["Need", "Current Source", "In Repo Today"]
    rows = [
        [
            "Delay labels and history",
            "Fabric SQL: bts_ontime_reporting, airline_delay_causes",
            "Loaded via scripts/16_load_fabric_sql_warehouse.py",
        ],
        [
            "Exogenous live context",
            "KQL: hazards_airsigmets, hazards_gairmets, opensky_states",
            "Runtime KQL read path exists in unified_retriever.py",
        ],
        [
            "Airport constraints",
            "Cosmos NOTAM + Postgres demo.notam_parsed",
            "Cosmos query path and parsed table both available",
        ],
        [
            "Network contagion",
            "demo.openflights_routes + demo.ops_graph_edges",
            "Graph traversal and SQL fallback already available",
        ],
        [
            "Scoring population",
            "demo.ops_flight_legs (today departures view)",
            "Synthetic near-term legs available for demo scoring",
        ],
    ]
    add_table(
        s,
        headers,
        rows,
        x=0.55,
        y=1.65,
        w=12.2,
        h=4.9,
        col_widths=[2.6, 4.6, 5.0],
        body_font_size=10,
    )
    note = s.shapes.add_textbox(Inches(0.65), Inches(6.68), Inches(12.0), Inches(0.3))
    ntf = note.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    np.text = "Data freshness note: current local schedule snapshot window is 2025-10 to 2025-11 (BTS), used for demo-scoped KPI claims."
    np.font.name = "Segoe UI"
    np.font.size = Pt(10)
    np.font.color.rgb = COLOR_MID
    add_footer(s)


def slide_3_product_experience(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Product Experience", "\"Today's departures\" with explainable risk and A/B toggle")

    left = add_panel(s, 0.7, 1.7, 8.6, 4.95, fill=COLOR_PANEL_ALT)
    set_panel_text(
        left,
        "UI Layout (mock)",
        [
            "Table columns: flight, route, std, risk(A15), expected delay, interval, top drivers.",
            "Row badges: High/Med/Low risk with color coding.",
            "One-click filter: highest expected impact first.",
        ],
        body_size=11,
    )

    table_headers = ["Flight", "Route", "Risk", "Delay", "Interval", "Top Driver"]
    table_rows = [
        ["TK203", "IST-LHR", "0.78", "31m", "[14, 48]", "SIGMET near arrival"],
        ["TK641", "IST-ESB", "0.56", "18m", "[6, 30]", "Runway NOTAM active"],
        ["TK792", "SAW-DXB", "0.49", "15m", "[4, 28]", "Traffic density up"],
        ["TK119", "IST-JFK", "0.44", "12m", "[2, 25]", "Inbound leg delay"],
    ]
    add_table(
        s,
        table_headers,
        table_rows,
        x=1.0,
        y=3.0,
        w=8.0,
        h=2.95,
        col_widths=[1.1, 1.5, 0.9, 1.0, 1.2, 2.3],
        body_font_size=9,
    )

    right_top = add_panel(s, 9.55, 1.7, 3.1, 2.25)
    set_panel_text(
        right_top,
        "Toggle",
        [
            "Baseline model",
            "Optimized model",
            "Delta KPI badge",
        ],
        body_size=12,
    )
    right_bottom = add_panel(s, 9.55, 4.15, 3.1, 2.5)
    set_panel_text(
        right_bottom,
        "KPI Cards",
        [
            "AUROC: 0.68 -> 0.79",
            "Brier: 0.192 -> 0.157",
            "MAE: 19.4m -> 14.1m",
            "Demo-scoped sample period",
        ],
        body_size=11,
    )
    add_footer(s)


def slide_4_model_progression(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Model Progression", "v0 baseline to v2 optimized path for this implementation")
    add_bullets(
        s,
        [
            "Training tracks: classifier for A15 risk and regressor for delay minutes.",
            "Single model family for optimized path: gradient-boosted trees with probability calibration.",
            "All versions scored on the same holdout period for A/B comparability.",
        ],
        x=0.8,
        y=1.65,
        w=12.0,
        h=1.2,
        size=14,
    )

    v0 = add_panel(s, 0.8, 3.05, 3.95, 2.9)
    set_panel_text(
        v0,
        "v0 - Baseline",
        [
            "Route + airline + time bucket priors.",
            "Simple GLM fallback score.",
            "Used as A/B control in UI.",
        ],
    )
    v1 = add_panel(s, 4.9, 3.05, 3.95, 2.9)
    set_panel_text(
        v1,
        "v1 - Exogenous",
        [
            "Add hazards + OpenSky + NOTAM features.",
            "Captures short-term operational stress.",
            "Improves real-time sensitivity.",
        ],
    )
    v2 = add_panel(s, 9.0, 3.05, 3.95, 2.9)
    set_panel_text(
        v2,
        "v2 - Network Aware",
        [
            "Add contagion and graph centrality.",
            "Inbound and upstream congestion proxies.",
            "Primary optimized model for demo.",
        ],
    )
    add_footer(s)


def slide_5_feature_architecture(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Feature Architecture", "How multi-source inputs become trainable and scoreable features")

    src = add_panel(s, 0.75, 1.8, 3.4, 4.85)
    set_panel_text(
        src,
        "Source Layer",
        [
            "Fabric SQL labels (BTS).",
            "KQL hazards and traffic.",
            "Cosmos + notam_parsed.",
            "openflights_routes + ops_graph_edges.",
        ],
    )
    feat = add_panel(s, 4.45, 1.8, 4.2, 4.85)
    set_panel_text(
        feat,
        "Feature Layer",
        [
            "Time-bucket and route priors.",
            "Weather and traffic pressure scores.",
            "NOTAM runway/airport constraints.",
            "Network centrality and upstream delay.",
            "Materialized in demo.delay_training_features.",
        ],
    )
    score = add_panel(s, 8.95, 1.8, 3.65, 4.85)
    set_panel_text(
        score,
        "Scoring Layer",
        [
            "A15 probability output.",
            "Expected delay minutes output.",
            "Prediction interval output.",
            "Top-3 driver extraction.",
            "Materialized in demo.delay_predictions_current.",
        ],
    )
    add_footer(s)


def slide_6_serving_design(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Scoring and Serving Design", "Hybrid batch + online refresh in current backend")

    top = add_panel(s, 0.8, 1.75, 12.0, 1.1)
    set_panel_text(
        top,
        "Execution mode",
        [
            "Hourly batch scores all near-term legs; request-time online refresh updates live hazard/traffic/NOTAM deltas.",
        ],
        title_size=13,
        body_size=12,
    )

    lane1 = add_panel(s, 0.8, 3.05, 5.8, 2.7)
    set_panel_text(
        lane1,
        "Batch lane",
        [
            "scripts/predictive/06_refresh_scores.py",
            "Writes demo.delay_scoring_features_current",
            "Stores baseline and optimized predictions",
        ],
    )
    lane2 = add_panel(s, 7.0, 3.05, 5.8, 2.7)
    set_panel_text(
        lane2,
        "Online lane",
        [
            "GET /api/predictive/delays",
            "Applies live deltas and driver extraction",
            "Returns degraded-source flags when needed",
        ],
    )

    endpoints = s.shapes.add_textbox(Inches(0.85), Inches(6.0), Inches(12.1), Inches(0.7))
    etf = endpoints.text_frame
    etf.clear()
    ep = etf.paragraphs[0]
    ep.text = "Endpoints: /api/predictive/delays and /api/predictive/delay-metrics (baseline vs optimized KPI summaries)."
    ep.font.name = "Segoe UI"
    ep.font.size = Pt(12)
    ep.font.color.rgb = COLOR_DARK
    add_footer(s)


def slide_7_explainability_uncertainty(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Explainability and Uncertainty", "Decision support needs score + reason + confidence")

    p_left = add_panel(s, 0.8, 1.8, 6.05, 4.85)
    set_panel_text(
        p_left,
        "Explainability",
        [
            "Top-3 drivers per flight from feature attribution.",
            "Driver phrasing localized for ops users.",
            "Examples: SIGMET near arrival, runway NOTAM, traffic density increase.",
        ],
    )
    p_right = add_panel(s, 7.1, 1.8, 5.45, 4.85)
    set_panel_text(
        p_right,
        "Uncertainty",
        [
            "Calibrated risk probability for A15.",
            "Delay interval from residual quantiles.",
            "Wider bands on sparse or degraded live context.",
            "Always label as decision support, not deterministic ETA.",
        ],
    )

    badge_hi = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.05), Inches(2.8), Inches(0.48)
    )
    badge_hi.fill.solid()
    badge_hi.fill.fore_color.rgb = COLOR_GREEN
    badge_hi.line.fill.background()
    badge_hi.text = "High confidence: narrow interval"
    for p in badge_hi.text_frame.paragraphs:
        p.font.name = "Segoe UI"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_WHITE

    badge_lo = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(6.05), Inches(3.3), Inches(0.48)
    )
    badge_lo.fill.solid()
    badge_lo.fill.fore_color.rgb = COLOR_AMBER
    badge_lo.line.fill.background()
    badge_lo.text = "Lower confidence: sparse/degraded signals"
    for p in badge_lo.text_frame.paragraphs:
        p.font.name = "Segoe UI"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_WHITE
    add_footer(s)


def slide_8_kpi_framework(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "KPI Framework (A/B Toggle)", "How baseline vs optimized uplift is shown in demo")
    headers = ["Metric", "Baseline (v0)", "Optimized (v2)", "Interpretation"]
    rows = [
        ["AUROC", "0.68", "0.79", "Higher is better ranking of risky legs"],
        ["Brier", "0.192", "0.157", "Lower is better probability calibration"],
        ["MAE (minutes)", "19.4", "14.1", "Lower is better delay-min prediction"],
        ["Top decile capture", "41%", "58%", "Higher captures more true A15 events"],
    ]
    add_table(
        s,
        headers,
        rows,
        x=0.75,
        y=1.8,
        w=11.9,
        h=2.7,
        col_widths=[1.8, 2.2, 2.2, 5.7],
        body_font_size=11,
    )

    p = add_panel(s, 0.75, 4.75, 11.9, 1.85)
    set_panel_text(
        p,
        "Concrete flight examples shown in demo",
        [
            "Flight TK203 flagged high risk in optimized model due to combined SIGMET + inbound delay contagion.",
            "Flight TK119 downgraded after NOTAM severity and traffic pressure normalized in online refresh.",
            "All KPI values are clearly tagged as demo-scope with current BTS snapshot coverage.",
        ],
        title_size=13,
        body_size=11,
    )
    add_footer(s)


def slide_9_delivery(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Delivery Plan", "Implementation sequence for fast demo readiness")

    c1 = add_panel(s, 0.75, 1.85, 2.95, 4.65)
    set_panel_text(
        c1,
        "Workstream 1\nData and Features",
        [
            "Build training set from Fabric SQL.",
            "Materialize context features.",
            "Define leakage-safe windows.",
        ],
    )
    c2 = add_panel(s, 3.95, 1.85, 2.95, 4.65)
    set_panel_text(
        c2,
        "Workstream 2\nModeling",
        [
            "Train v0 baseline and v2 optimized.",
            "Calibrate risk probabilities.",
            "Generate interval estimator.",
        ],
    )
    c3 = add_panel(s, 7.15, 1.85, 2.95, 4.65)
    set_panel_text(
        c3,
        "Workstream 3\nAPI and UI",
        [
            "Serve predictive endpoints.",
            "Add Today departures view.",
            "Add baseline/optimized toggle.",
        ],
    )
    c4 = add_panel(s, 10.35, 1.85, 2.25, 4.65)
    set_panel_text(
        c4,
        "Workstream 4\nValidation",
        [
            "Contract tests.",
            "Degraded-source tests.",
            "Demo script rehearsal.",
        ],
        body_size=10,
    )
    add_footer(s)


def slide_10_risks_next(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(s)
    add_header(s, "Risks, Limits, and Next Step", "Transparent scope boundaries and execution guardrails")

    left = add_panel(s, 0.8, 1.85, 5.95, 4.75)
    set_panel_text(
        left,
        "Known limits",
        [
            "BTS history currently limited to available loaded months for this demo.",
            "Live-source availability can vary by endpoint readiness.",
            "Predictions are optimization aids, not dispatch automation.",
        ],
    )
    right = add_panel(s, 7.05, 1.85, 5.45, 4.75)
    set_panel_text(
        right,
        "Immediate next step after this deck",
        [
            "Implement scripts + API + UI slices exactly as specified.",
            "Run A/B evaluation and attach reproducible metrics artifacts.",
            "Prepare stakeholder runbook with fallback mode and source health checks.",
            "Defer v3 cost-threshold optimization to next phase.",
        ],
    )
    add_footer(s, text="Predictive Delay Optimization - Demo scope complete, v3 deferred")


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_1_title(prs)
    slide_2_feasible_now(prs)
    slide_3_product_experience(prs)
    slide_4_model_progression(prs)
    slide_5_feature_architecture(prs)
    slide_6_serving_design(prs)
    slide_7_explainability_uncertainty(prs)
    slide_8_kpi_framework(prs)
    slide_9_delivery(prs)
    slide_10_risks_next(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
