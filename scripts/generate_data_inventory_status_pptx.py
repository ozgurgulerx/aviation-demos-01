#!/usr/bin/env python3
"""
Generate a concise data inventory + Fabric gap status deck.
"""

from __future__ import annotations

from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/THY_Data_Inventory_Storage_Status_Fabric_Gap.pptx")

COLOR_RED = RGBColor(198, 12, 48)
COLOR_DARK = RGBColor(30, 35, 45)
COLOR_LIGHT = RGBColor(245, 247, 250)
COLOR_MID = RGBColor(99, 110, 126)
COLOR_WHITE = RGBColor(255, 255, 255)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT
    bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
    top.fill.solid()
    top.fill.fore_color.rgb = COLOR_RED
    top.line.fill.background()


def add_header(slide, title: str, subtitle: str = ""):
    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(0.8)).text_frame
    t.clear()
    p = t.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Segoe UI"
    r.font.bold = True
    r.font.size = Pt(30)
    r.font.color.rgb = COLOR_DARK

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.72), Inches(1.05), Inches(12.0), Inches(0.5)).text_frame
        s.clear()
        sp = s.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Segoe UI"
        sr.font.size = Pt(14)
        sr.font.color.rgb = COLOR_MID


def add_footer(slide, text="Turkish Airlines Demo - Data Inventory Status"):
    f = slide.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(12.0), Inches(0.25)).text_frame
    f.clear()
    p = f.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(10)
    r.font.color.rgb = COLOR_MID


def add_bullets(slide, bullets: List[str], x=0.8, y=1.7, w=12.0, h=5.4, size=18):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = COLOR_DARK


def add_table_slide(slide, title: str, subtitle: str, rows: List[List[str]], part: str):
    add_bg(slide)
    add_header(slide, title, f"{subtitle} - {part}")
    cols = ["Dataset", "Source Path", "Current Store", "Count"]
    table = slide.shapes.add_table(len(rows) + 1, len(cols), Inches(0.45), Inches(1.6), Inches(12.45), Inches(5.3)).table

    widths = [2.5, 4.6, 3.1, 1.7]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    for c, h in enumerate(cols):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_RED
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE

    for r_i, row in enumerate(rows, start=1):
        for c_i, value in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = value
            if r_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(235, 239, 245)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(10)
                p.font.color.rgb = COLOR_DARK

    add_footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s1)
    add_header(
        s1,
        "Data Inventory and Storage Status",
        "What is loaded now vs what is Fabric-native",
    )
    add_bullets(
        s1,
        [
            "This deck lists all datasets used in the pilot brief demo and their current datastore placement.",
            "Current state is hybrid: Azure PostgreSQL + AI Search + Table Storage + ADX/Kusto external tables.",
            "Reason you see more PostgreSQL today: Fabric/Cosmos provisioning and API constraints in this tenant/session.",
            "Fabric-forward mapping is included in the final slides with concrete next actions.",
        ],
        size=18,
    )
    add_footer(s1)

    all_rows = [
        ["ASRS reports", "data/c1-asrs/processed/asrs_records.jsonl", "PostgreSQL public.asrs_reports", "150,257"],
        ["Ops vector docs", "data/vector_docs/ops_narratives_docs.jsonl", "AI Search idx_ops_narratives", "3,000"],
        ["Regulatory vector docs", "data/vector_docs/regulatory_docs.jsonl", "AI Search idx_regulatory", "55"],
        ["Airport vector docs", "data/vector_docs/airport_ops_docs.jsonl", "AI Search idx_airport_ops_docs", "2,000"],
        ["OurAirports airports", "data/g-ourairports_recent/airports_*.csv", "PostgreSQL demo.ourairports_airports", "84,578"],
        ["OurAirports runways", "data/g-ourairports_recent/runways_*.csv", "PostgreSQL demo.ourairports_runways", "47,602"],
        ["OurAirports navaids", "data/g-ourairports_recent/navaids_*.csv", "PostgreSQL demo.ourairports_navaids", "11,010"],
        ["OurAirports frequencies", "data/g-ourairports_recent/airport-frequencies_*.csv", "PostgreSQL demo.ourairports_frequencies", "30,211"],
        ["OpenFlights routes", "data/f-openflights/raw/routes_*.dat", "PostgreSQL demo.openflights_routes", "67,663"],
        ["OpenFlights airports", "data/f-openflights/raw/airports_*.dat", "PostgreSQL demo.openflights_airports", "7,698"],
        ["OpenFlights airlines", "data/f-openflights/raw/airlines_*.dat", "PostgreSQL demo.openflights_airlines", "6,162"],
        ["OpenSky raw JSON", "data/e-opensky_recent/*.json", "PostgreSQL demo.opensky_raw", "56"],
        ["NOTAM Istanbul", "data/h-notam_recent/.../search_location_istanbul.jsonl", "PostgreSQL demo.notam_raw", "30"],
        ["Hazards AIRMET/SIGMET", "data/i-aviationweather_hazards_recent/airsigmets*.gz", "PostgreSQL demo.hazards_airsigmets", "3"],
        ["Hazards G-AIRMET", "data/i-aviationweather_hazards_recent/gairmets*.gz", "PostgreSQL demo.hazards_gairmets", "208"],
        ["Hazards PIREP raw", "data/i-aviationweather_hazards_recent/aircraftreports*.gz", "PostgreSQL demo.hazards_aireps_raw", "2,039"],
        ["Synthetic flight legs", "data/j-synthetic_ops_overlay/.../ops_flight_legs.csv", "PostgreSQL demo.ops_flight_legs", "173"],
        ["Synthetic milestones", "data/j-synthetic_ops_overlay/.../ops_turnaround_milestones.csv", "PostgreSQL demo.ops_turnaround_milestones", "1,384"],
        ["Synthetic baggage events", "data/j-synthetic_ops_overlay/.../ops_baggage_events.csv", "PostgreSQL demo.ops_baggage_events", "682"],
        ["Synthetic crew rosters", "data/j-synthetic_ops_overlay/.../ops_crew_rosters.csv", "PostgreSQL demo.ops_crew_rosters", "519"],
        ["Synthetic MEL/techlog", "data/j-synthetic_ops_overlay/.../ops_mel_techlog_events.csv", "PostgreSQL demo.ops_mel_techlog_events", "41"],
        ["Synthetic graph edges", "data/j-synthetic_ops_overlay/.../ops_graph_edges.csv", "PostgreSQL demo.ops_graph_edges", "519"],
        ["OpenSky state vectors", "data/e-opensky_recent/opensky_states_all_*.json", "Table Storage openskystates", "10,866"],
        ["NOTAM NoSQL copy", "data/h-notam_recent/.../search_location_istanbul.jsonl", "Table Storage notamdocs", "31"],
        ["Graph edge NoSQL copy", "data/j-synthetic_ops_overlay/.../ops_graph_edges.csv", "Table Storage opsgraphedges", "519"],
        ["Kusto external table files", "Blob kustoingest20260219/*.csv", "ADX/Kusto external tables", "7 blobs"],
    ]

    chunks = [all_rows[i : i + 9] for i in range(0, len(all_rows), 9)]
    for i, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_table_slide(slide, "Current Data Placement", "Loaded datasets", chunk, f"Part {i}/{len(chunks)}")

    # Fabric gap slide
    s_gap = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s_gap)
    add_header(s_gap, "Why You See More PostgreSQL Than Fabric", "Current blockers and rationale")
    add_bullets(
        s_gap,
        [
            "1) Cosmos DB account provisioning failed due regional capacity/quota in this subscription at runtime.",
            "2) Fabric/Kusto data-plane endpoint connectivity from this runner timed out for direct ingest calls.",
            "3) To avoid blocking demo readiness, ingestion path shifted to deterministic stores that were provisionable immediately.",
            "4) Result: retrieval works now, but Fabric-native storage coverage is not yet at target breadth.",
        ],
        size=17,
    )
    add_footer(s_gap)

    # Fabric plan slide
    s_plan = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s_plan)
    add_header(s_plan, "Fabric-First Target Mapping", "What should move next")
    plan_rows = [
        ["OpenSky + hazards streams", "Fabric Eventstream + Eventhouse", "Partially ready", "Create stream connectors and ingestion jobs"],
        ["Airport/reference dimensions", "Fabric Warehouse/Lakehouse", "Ready", "COPY from current CSV/SQL staging"],
        ["Graph edges + relationships", "Fabric Graph (preview)", "Ready to migrate", "Load from demo.ops_graph_edges"],
        ["NOTAM document corpus", "Fabric Lakehouse + OneLake shortcuts", "Ready", "Move JSONL + index in Fabric-connected search"],
        ["Ops synthetic overlay", "Fabric Warehouse", "Ready", "Bulk load from PostgreSQL demo schema"],
        ["RAG vector indexes", "Keep in Azure AI Search", "Active", "Optional Fabric IQ bridging, keep hybrid retrieval"],
    ]
    add_table_slide(s_plan, "Fabric-First Target Mapping", "Proposed migration sequencing", plan_rows, "Execution View")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()

